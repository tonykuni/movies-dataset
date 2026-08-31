#!/usr/bin/env python3
"""VIA Financial Document Intelligence & Four-Bullet Summarizer Engine.

單檔、離線優先、選配依賴可降級。功能：文件攝取、Layout/階層辨識、文字修復、
表圖標註、台股官方主檔解析、金融硬數據 Regex、抽取式或 GGUF 摘要、格式驗證與輸出。
"""

from __future__ import annotations
# ===== [VIA:ACCEL-BRIDGE:v0100] SuperAccel 加速器橋(批102 全樹導入令;graceful 零行為變更) =====
try:
    import sys as _sa_sys
    from pathlib import Path as _sa_Path
    _sa_p = _sa_Path(__file__).resolve()
    while _sa_p.parent != _sa_p:
        if (_sa_p / "supportive modules" / "VIA_SuperAccel_Module.py").exists():
            _sa_sys.path.insert(0, str(_sa_p / "supportive modules"))
            break
        _sa_p = _sa_p.parent
    import VIA_SuperAccel_Module as VIA_ACCEL  # noqa: N816
except Exception:
    VIA_ACCEL = None  # graceful:加速器缺席零影響
# ===== [VIA:ACCEL-BRIDGE:END] =====

# =============================================================================
# def 00 PARAMETERS / SSOT
# =============================================================================

import argparse
import csv
import hashlib
import html
import importlib.util
import json
import logging
import math
import os
import re
import statistics
import sys
import unicodedata
from collections import Counter
from dataclasses import asdict, dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, Optional, Protocol, Sequence

ENGINE_NAME = "VIA Financial Document Intelligence & Four-Bullet Summarizer"
ENGINE_VERSION = "1.0.0"
DEFAULT_ENCODING = "utf-8-sig"
DEFAULT_BACKEND = "auto"
DEFAULT_LANGUAGE = "zh-en"
DEFAULT_MAX_INPUT_CHARS = 120_000
DEFAULT_MAX_SENTENCE_CHARS = 320
DEFAULT_THREADS = max(1, (os.cpu_count() or 4) // 2)
DEFAULT_CONTEXT = 4096
DEFAULT_MAX_TOKENS = 720
DEFAULT_TEMPERATURE = 0.1
DEFAULT_TOP_P = 0.9
DEFAULT_FOOTER_RATIO = 0.92
DEFAULT_HEADER_RATIO = 0.04
DEFAULT_SIDE_RATIO = 0.08
DEFAULT_MIN_TEXT_CHARS = 2
MISSING_VALUE = "未提及"
REVIEW_REQUIRED = "REVIEW_REQUIRED"

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".xlsm", ".pptx", ".odt", ".ods", ".odp",
    ".html", ".htm", ".md", ".markdown", ".txt", ".json", ".csv", ".tsv",
}

TABLE_CAPTION_RE = re.compile(r"^(?:表|Table)\s*([A-Z]?\d+(?:[-.]\d+)*)\s*[：:.、-]?\s*(.*)$", re.I)
FIGURE_CAPTION_RE = re.compile(r"^(?:圖|Figure|Fig\.?|Chart)\s*([A-Z]?\d+(?:[-.]\d+)*)\s*[：:.、-]?\s*(.*)$", re.I)
SOURCE_RE = re.compile(r"^(?:資料來源|來源|Source|Sources)\s*[：:]\s*(.+)$", re.I)
DISCLAIMER_RE = re.compile(r"免責|Disclaimer|重要聲明|請詳閱.*尾頁|copyright|版權所有", re.I)
TICKER_RE = re.compile(r"(?<!\d)(?P<ticker>\d{4}[A-Z]?)(?:\s*(?:\.|\s)(?P<suffix>TWO|TW|TT))?(?!\d)", re.I)
RATING_RE = re.compile(r"(?:評等|Rating)\s*[：:]?\s*(買進|增持|中立|持有|減持|賣出|Buy|Outperform|Neutral|Hold|Underperform|Sell)", re.I)
TARGET_RE = re.compile(r"(?:目標價|合理價|推估價|TP|Target\s*Price)\s*(?:為|上看|調升至|調降至|：|:)?\s*(?:NT\$|TWD|新台幣)?\s*(?P<value>\d+(?:,\d{3})*(?:\.\d+)?)\s*(?P<unit>元|塊|美元|USD|NTD|TWD)?", re.I)
CURRENT_RE = re.compile(r"(?:目前股價|現價|收盤價|Current\s*Price)\s*(?:為|：|:)?\s*(?:NT\$|TWD)?\s*(?P<value>\d+(?:,\d{3})*(?:\.\d+)?)", re.I)
UPSIDE_RE = re.compile(r"(?P<value>[-+]?\d+(?:\.\d+)?)\s*%\s*(?:的)?(?:上漲空間|漲幅|潛在報酬|upside)?", re.I)
VALUATION_RE = re.compile(r"(?P<method>(?:基於|採用|給予|以|根據)?[^。；;\n]{0,55}?(?:本益比|股價淨值比|P/?E|PER|PBR|DCF|EV/?EBITDA|SOTP)[^。；;\n]{0,80})", re.I)
FINANCIAL_RE = re.compile(r"(?:營收|毛利率|營益率|淨利|EPS|每股盈餘|資本支出|CAPEX)[^。！？!?\n]{0,180}", re.I)
CATALYST_RE = re.compile(r"(?:催化劑|法說會|量產|出貨|漲價|調升|上修|下修|新產品|訂單|產能|市佔)[^。！？!?\n]{0,180}", re.I)
RISK_RE = re.compile(r"(?:風險|下行|不確定|庫存|匯率|地緣政治|需求疲弱|遞延|競爭|Risk)[^。！？!?\n]{0,180}", re.I)

SUMMARY_LABELS = ("評價與目標", "核心驅動", "營運與財務亮點", "風險與展望")


# =============================================================================
# def 01 DATA CONTRACTS
# =============================================================================

@dataclass
class DocumentBlock:
    page_num: int
    block_id: str
    content_type: str
    sub_type: str
    text_content: str
    font_size: Optional[float] = None
    is_bold: bool = False
    bbox: Optional[list[float]] = None
    table_no: Optional[str] = None
    figure_no: Optional[str] = None
    data_source: Optional[str] = None
    confidence: float = 1.0
    excluded: bool = False
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class SecurityIdentity:
    ticker: str = ""
    short_name: str = ""
    market: str = ""
    yfinance_ticker: str = ""
    source: str = ""
    status: str = REVIEW_REQUIRED


@dataclass
class InvestmentMetrics:
    rating: str = MISSING_VALUE
    target_price: str = MISSING_VALUE
    current_price: str = MISSING_VALUE
    upside: str = MISSING_VALUE
    valuation_method: str = MISSING_VALUE
    valuation_reason: str = MISSING_VALUE
    evidence: dict[str, str] = field(default_factory=dict)


@dataclass
class EngineResult:
    run_id: str
    source_file: str
    source_sha256: str
    engine: str
    version: str
    created_at_utc: str
    parser: str
    summarizer: str
    security: SecurityIdentity
    header: str
    metrics: InvestmentMetrics
    summary_bullets: list[str]
    final_text: str
    blocks: list[DocumentBlock]
    quality_gate: str
    findings: list[str]
    tool_registry: dict[str, Any]


class SummarizerBackend(Protocol):
    name: str
    def summarize(self, context: "SummaryContext") -> list[str]: ...


@dataclass
class SummaryContext:
    header: str
    metrics: InvestmentMetrics
    clean_text: str
    sentences: list[str]


# =============================================================================
# def 02 TOOL REGISTRY / CAPABILITY DISCOVERY
# =============================================================================

def def_has_module(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def def_build_tool_registry() -> dict[str, Any]:
    candidates = {
        "layout": {"pymupdf": "fitz", "pdfplumber": "pdfplumber", "surya": "surya", "layoutparser": "layoutparser"},
        "tables": {"camelot": "camelot", "tabula": "tabula", "pandas": "pandas"},
        "ocr": {"rapidocr": "rapidocr_onnxruntime", "pytesseract": "pytesseract", "paddleocr": "paddleocr"},
        "office": {"python-docx": "docx", "openpyxl": "openpyxl", "python-pptx": "pptx", "odfpy": "odf"},
        "html": {"beautifulsoup4": "bs4", "lxml": "lxml", "markdown-it": "markdown_it"},
        "extractive": {"sumy": "sumy", "jieba": "jieba", "snownlp": "snownlp", "pytextrank": "pytextrank", "summa": "summa"},
        "generative": {"llama-cpp-python": "llama_cpp", "transformers": "transformers", "langchain": "langchain", "ctransformers": "ctransformers", "onnxruntime": "onnxruntime"},
        "storage": {"pyarrow": "pyarrow", "duckdb": "duckdb"},
    }
    return {
        category: {display: {"module": module, "available": def_has_module(module)} for display, module in libs.items()}
        for category, libs in candidates.items()
    }


# =============================================================================
# def 03 NORMALIZATION / REPAIR
# =============================================================================

def def_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def def_normalize_text(text: Any) -> str:
    value = "" if text is None else str(text)
    value = unicodedata.normalize("NFKC", value).replace("\u00ad", "").replace("\u200b", "")
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    value = re.sub(r"(?<=\w)-\n(?=\w)", "", value)
    value = re.sub(r"[ \t\f\v]+", " ", value)
    value = re.sub(r" *\n *", "\n", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def def_repair_mojibake(text: str) -> str:
    if not text or "�" not in text:
        return text
    if def_has_module("ftfy"):
        from ftfy import fix_text
        return fix_text(text)
    return text


def def_split_sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[。！？!?；;])\s*|\n+", def_normalize_text(text))
    return [p.strip(" *•-\t") for p in parts if len(p.strip()) >= DEFAULT_MIN_TEXT_CHARS]


def def_dedupe_sentences(sentences: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    output: list[str] = []
    for sentence in sentences:
        key = re.sub(r"\W+", "", sentence).lower()
        if key and key not in seen:
            seen.add(key)
            output.append(sentence)
    return output


# =============================================================================
# def 04 DOCUMENT PARSERS
# =============================================================================

def def_classify_text(text: str, size: Optional[float], base: Optional[float], max_size: Optional[float], bold: bool) -> tuple[str, str]:
    if TABLE_CAPTION_RE.match(text): return "文", "表頭(含表號)"
    if FIGURE_CAPTION_RE.match(text): return "文", "圖頭(含圖號)"
    if SOURCE_RE.match(text): return "文", "表/圖資料來源"
    if DISCLAIMER_RE.search(text): return "文", "免責聲明"
    if size is None or base is None: return "文", "本文"
    if max_size and size >= max(max_size * 0.92, base * 1.45): return "文", "大標題"
    if bold and size >= base * 1.16: return "文", "中標題"
    if bold and size >= base * 0.96: return "文", "小標題"
    if size < base * 0.75: return "文", "小文字"
    return "文", "本文"


def def_make_plain_blocks(text: str, source: str = "plain") -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    for idx, line in enumerate((x.strip() for x in text.splitlines()), 1):
        if not line: continue
        content_type, sub_type = def_classify_text(line, None, None, None, False)
        blocks.append(DocumentBlock(1, f"P1-B{idx:04d}", content_type, sub_type, line, metadata={"parser": source}))
    return blocks


def def_parse_pdf(path: Path) -> tuple[list[DocumentBlock], str]:
    if not def_has_module("fitz"):
        raise RuntimeError("PDF 解析需要 PyMuPDF：pip install pymupdf")
    import fitz
    doc = fitz.open(path)
    spans: list[dict[str, Any]] = []
    for page_index, page in enumerate(doc, 1):
        data = page.get_text("dict", sort=True)
        for block in data.get("blocks", []):
            if block.get("type") != 0: continue
            for line in block.get("lines", []):
                line_spans = [s for s in line.get("spans", []) if s.get("text", "").strip()]
                if not line_spans: continue
                text = def_normalize_text("".join(s.get("text", "") for s in line_spans))
                sizes = [float(s.get("size", 0)) for s in line_spans if s.get("size")]
                fonts = " ".join(str(s.get("font", "")) for s in line_spans)
                spans.append({"page": page_index, "text": text, "size": max(sizes) if sizes else None,
                              "bold": bool(re.search(r"bold|black|heavy|demi", fonts, re.I)),
                              "bbox": list(line.get("bbox", block.get("bbox", []))),
                              "height": float(page.rect.height), "width": float(page.rect.width)})
    sizes = [round(x["size"], 1) for x in spans if x["size"]]
    base = Counter(sizes).most_common(1)[0][0] if sizes else None
    maximum = max(sizes) if sizes else None
    blocks: list[DocumentBlock] = []
    page_counts: Counter[int] = Counter()
    for item in spans:
        page_counts[item["page"]] += 1
        bbox = item["bbox"]
        excluded = False
        subtype_override = ""
        if bbox and bbox[3] >= item["height"] * DEFAULT_FOOTER_RATIO:
            excluded, subtype_override = True, "頁尾"
        elif bbox and bbox[1] <= item["height"] * DEFAULT_HEADER_RATIO and item["page"] > 1:
            excluded, subtype_override = True, "頁首"
        ctype, subtype = def_classify_text(item["text"], item["size"], base, maximum, item["bold"])
        subtype = subtype_override or subtype
        table_match, figure_match, source_match = TABLE_CAPTION_RE.match(item["text"]), FIGURE_CAPTION_RE.match(item["text"]), SOURCE_RE.match(item["text"])
        blocks.append(DocumentBlock(item["page"], f"P{item['page']}-B{page_counts[item['page']]:04d}", ctype, subtype,
                                    item["text"], item["size"], item["bold"], bbox,
                                    table_no=table_match.group(1) if table_match else None,
                                    figure_no=figure_match.group(1) if figure_match else None,
                                    data_source=source_match.group(1) if source_match else None,
                                    excluded=excluded or subtype in {"免責聲明", "小文字"}, metadata={"parser": "pymupdf"}))
    # pdfplumber is optional and enriches actual tables without becoming a hard dependency.
    if def_has_module("pdfplumber"):
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page_no, page in enumerate(pdf.pages, 1):
                for table_no, table in enumerate(page.extract_tables() or [], 1):
                    clean = [[def_normalize_text(cell) for cell in row] for row in table if row]
                    blocks.append(DocumentBlock(page_no, f"P{page_no}-T{table_no:03d}", "表", "表格資料", json.dumps(clean, ensure_ascii=False), metadata={"parser": "pdfplumber"}))
    doc.close()
    return blocks, "pymupdf+pdfplumber" if def_has_module("pdfplumber") else "pymupdf"


def def_parse_docx(path: Path) -> tuple[list[DocumentBlock], str]:
    if not def_has_module("docx"): raise RuntimeError("DOCX 解析需要 python-docx")
    from docx import Document
    document = Document(path)
    blocks: list[DocumentBlock] = []
    idx = 0
    for para in document.paragraphs:
        text = def_normalize_text(para.text)
        if not text: continue
        idx += 1
        style = (para.style.name or "").lower()
        subtype = "大標題" if "title" in style or "heading 1" in style else "中標題" if "heading 2" in style else "小標題" if "heading" in style else def_classify_text(text, None, None, None, False)[1]
        blocks.append(DocumentBlock(1, f"P1-B{idx:04d}", "文", subtype, text, metadata={"style": style, "parser": "python-docx"}))
    for t_idx, table in enumerate(document.tables, 1):
        rows = [[def_normalize_text(cell.text) for cell in row.cells] for row in table.rows]
        blocks.append(DocumentBlock(1, f"P1-T{t_idx:03d}", "表", "表格資料", json.dumps(rows, ensure_ascii=False), metadata={"parser": "python-docx"}))
    return blocks, "python-docx"


def def_parse_xlsx(path: Path) -> tuple[list[DocumentBlock], str]:
    if not def_has_module("openpyxl"): raise RuntimeError("XLSX 解析需要 openpyxl")
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    blocks = []
    for idx, ws in enumerate(wb.worksheets, 1):
        rows = [["" if v is None else str(v) for v in row] for row in ws.iter_rows(values_only=True)]
        blocks.append(DocumentBlock(idx, f"S{idx}-T001", "表", "工作表", json.dumps(rows, ensure_ascii=False), metadata={"sheet": ws.title, "parser": "openpyxl"}))
    wb.close()
    return blocks, "openpyxl"


def def_parse_pptx(path: Path) -> tuple[list[DocumentBlock], str]:
    if not def_has_module("pptx"): raise RuntimeError("PPTX 解析需要 python-pptx")
    from pptx import Presentation
    prs = Presentation(path)
    blocks = []
    for page, slide in enumerate(prs.slides, 1):
        idx = 0
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not def_normalize_text(shape.text): continue
            idx += 1
            text = def_normalize_text(shape.text)
            subtype = "大標題" if getattr(shape, "is_placeholder", False) and "TITLE" in str(getattr(shape, "placeholder_format", "")) else "本文"
            blocks.append(DocumentBlock(page, f"P{page}-B{idx:04d}", "文", subtype, text, metadata={"parser": "python-pptx"}))
    return blocks, "python-pptx"


def def_parse_html(path: Path) -> tuple[list[DocumentBlock], str]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    if def_has_module("bs4"):
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        for node in soup(["script", "style", "noscript"]): node.decompose()
        blocks = []
        for idx, node in enumerate(soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "table", "figcaption"]), 1):
            text = def_normalize_text(node.get_text(" ", strip=True))
            if not text: continue
            subtype = {"h1": "大標題", "h2": "中標題", "h3": "小標題", "h4": "小標題", "table": "表格資料", "figcaption": "圖頭(含圖號)"}.get(node.name, "本文")
            blocks.append(DocumentBlock(1, f"P1-B{idx:04d}", "表" if node.name == "table" else "文", subtype, text, metadata={"parser": "beautifulsoup4"}))
        return blocks, "beautifulsoup4"
    text = re.sub(r"<[^>]+>", "\n", html.unescape(raw))
    return def_make_plain_blocks(text, "html-regex-fallback"), "html-regex-fallback"


def def_parse_document(path: Path) -> tuple[list[DocumentBlock], str]:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS: raise ValueError(f"不支援的格式：{suffix}")
    if suffix == ".pdf": return def_parse_pdf(path)
    if suffix == ".docx": return def_parse_docx(path)
    if suffix in {".xlsx", ".xlsm"}: return def_parse_xlsx(path)
    if suffix == ".pptx": return def_parse_pptx(path)
    if suffix in {".html", ".htm"}: return def_parse_html(path)
    if suffix == ".json":
        data = json.loads(path.read_text(encoding="utf-8-sig"))
        return def_make_plain_blocks(json.dumps(data, ensure_ascii=False, indent=2), "json"), "json"
    if suffix in {".csv", ".tsv"}:
        delimiter = "\t" if suffix == ".tsv" else ","
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
            rows = list(csv.reader(handle, delimiter=delimiter))
        return [DocumentBlock(1, "P1-T001", "表", "表格資料", json.dumps(rows, ensure_ascii=False), metadata={"parser": "csv"})], "csv"
    if suffix in {".odt", ".ods", ".odp"}: raise RuntimeError("ODF 解析請先安裝 odfpy；此版本建議匯出為 DOCX/XLSX/PPTX")
    text = path.read_text(encoding="utf-8-sig", errors="replace")
    return def_make_plain_blocks(text, "plain-text"), "plain-text"


# =============================================================================
# def 05 SECURITY MASTER / HEADER
# =============================================================================

def def_load_security_master(path: Optional[Path]) -> dict[str, SecurityIdentity]:
    if not path: return {}
    if not path.exists(): raise FileNotFoundError(path)
    rows: list[dict[str, Any]]
    if path.suffix.lower() == ".json":
        obj = json.loads(path.read_text(encoding="utf-8-sig")); rows = obj if isinstance(obj, list) else obj.get("data", [])
    else:
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle: rows = list(csv.DictReader(handle))
    output = {}
    for row in rows:
        ticker = str(row.get("Ticker") or row.get("ticker") or row.get("代碼") or row.get("證券代號") or "").strip().upper()
        name = str(row.get("Name") or row.get("name") or row.get("中文簡稱") or row.get("證券名稱") or "").strip()
        market = str(row.get("Market") or row.get("market") or row.get("市場") or "").strip().upper()
        source = str(row.get("Source") or row.get("source") or "LOCAL_SECURITY_MASTER").strip()
        if ticker: output[ticker] = def_make_security_identity(ticker, name, market, source)
    return output


def def_make_security_identity(ticker: str, name: str, market: str, source: str) -> SecurityIdentity:
    normalized_market = {"上市": "TWSE", "TW": "TWSE", "上櫃": "TPEX", "OTC": "TPEX", "TWO": "TPEX"}.get(market.upper(), market.upper())
    yf = f"{ticker}.TW" if normalized_market == "TWSE" else f"{ticker}.TWO" if normalized_market == "TPEX" else ""
    status = "VERIFIED" if ticker and name and yf else REVIEW_REQUIRED
    return SecurityIdentity(ticker, name, normalized_market, yf, source, status)


def def_resolve_security(text: str, master: dict[str, SecurityIdentity], ticker_override: str = "") -> SecurityIdentity:
    ticker = ticker_override.strip().upper()
    if not ticker and master:
        # Official/local master is authoritative. This prevents years (2026) and
        # model parameters (2048) in long documents from becoming false tickers.
        positions = [(text.find(code), code) for code in master if re.search(rf"(?<!\d){re.escape(code)}(?!\d)", text)]
        if positions:
            ticker = min(positions, key=lambda item: item[0])[1]
    if not ticker and not master:
        explicit = re.search(r"(?<!\d)(?P<ticker>\d{4}[A-Z]?)\s*(?:\.|\s)(?:TW|TWO|TT)(?![A-Z0-9])", text, re.I)
        ticker = explicit.group("ticker").upper() if explicit else ""
    if ticker in master: return master[ticker]
    return SecurityIdentity(ticker=ticker, source="UNRESOLVED", status=REVIEW_REQUIRED)


def def_select_raw_title(blocks: Sequence[DocumentBlock]) -> str:
    active = [b for b in blocks if not b.excluded and b.content_type == "文" and len(b.text_content) >= 4]
    titles = [b for b in active if b.sub_type == "大標題"]
    if titles:
        return max(titles, key=lambda b: ((b.font_size or 0), b.is_bold, -b.page_num, min(len(b.text_content), 120))).text_content
    if not active: return "未識別標題"
    # Plain text has no geometry; document order is safer than selecting the longest line.
    return active[0].text_content


def def_clean_title(raw_title: str, security: SecurityIdentity, aliases: Sequence[str] = ()) -> str:
    title = raw_title
    keywords = [security.short_name, *aliases, "股份有限公司", "股份", "個股報告", "Company Update", "Equity Research"]
    for keyword in sorted((x for x in keywords if x), key=len, reverse=True): title = re.sub(re.escape(keyword), "", title, flags=re.I)
    if security.ticker:
        title = re.sub(rf"(?<!\d){re.escape(security.ticker)}(?:\s*(?:\.|\s)(?:TWO|TW|TT))?(?!\d)", "", title, flags=re.I)
    title = re.sub(r"[\[\]【】()（）]+", " ", title)
    title = re.sub(r"^[\s:：|｜,，.。;；—–_-]+|[\s:：|｜—–_-]+$", "", title)
    title = re.sub(r"\s+", " ", title).strip()
    return title or "未識別標題"


def def_build_header(title: str, security: SecurityIdentity) -> str:
    if security.status == "VERIFIED": return f"{security.short_name}({security.yfinance_ticker}): {title}"
    ticker = security.ticker or "UNKNOWN"
    return f"待確認({ticker}): {title}"


# =============================================================================
# def 06 REGEX HARD DATA
# =============================================================================

def def_first_group(pattern: re.Pattern[str], text: str, group: str = "value") -> tuple[str, str]:
    match = pattern.search(text)
    resolved_group: Any = int(group) if isinstance(group, str) and group.isdigit() else group
    return (match.group(resolved_group).strip(), match.group(0).strip()) if match else (MISSING_VALUE, "")


def def_extract_metrics(text: str) -> InvestmentMetrics:
    rating, rating_ev = def_first_group(RATING_RE, text, "1")
    target, target_ev = def_first_group(TARGET_RE, text)
    current, current_ev = def_first_group(CURRENT_RE, text)
    upside, upside_ev = def_first_group(UPSIDE_RE, text)
    valuation, valuation_ev = def_first_group(VALUATION_RE, text, "method")
    if target != MISSING_VALUE:
        target = target.replace(",", "") + (TARGET_RE.search(text).group("unit") or "元")
    if current != MISSING_VALUE: current = current.replace(",", "") + "元"
    if upside != MISSING_VALUE: upside = f"{upside}%"
    elif target != MISSING_VALUE and current != MISSING_VALUE:
        try:
            tp, cp = float(re.findall(r"\d+(?:\.\d+)?", target)[0]), float(re.findall(r"\d+(?:\.\d+)?", current)[0])
            if cp > 0: upside = f"{((tp / cp) - 1) * 100:.1f}%（由目標價與現價計算）"
        except (ValueError, IndexError): pass
    reason = MISSING_VALUE
    reason_match = re.search(r"(?:評價理由|主要基於|理由)[：:]?\s*([^。\n]{5,180})", text)
    if reason_match: reason = reason_match.group(1).strip()
    evidence = {k: v for k, v in {"rating": rating_ev, "target_price": target_ev, "current_price": current_ev,
                                   "upside": upside_ev, "valuation_method": valuation_ev}.items() if v}
    return InvestmentMetrics(rating, target, current, upside, valuation, reason, evidence)


# =============================================================================
# def 07 SUMMARIZERS
# =============================================================================

def def_sentence_score(sentence: str, keywords: Sequence[str], category_pattern: re.Pattern[str]) -> float:
    score = 2.5 if category_pattern.search(sentence) else 0.0
    score += sum(0.35 for keyword in keywords if keyword and keyword.lower() in sentence.lower())
    score += min(len(re.findall(r"\d+(?:\.\d+)?%?", sentence)), 4) * 0.45
    score += 0.5 if 18 <= len(sentence) <= DEFAULT_MAX_SENTENCE_CHARS else 0
    return score


def def_pick_sentence(sentences: Sequence[str], pattern: re.Pattern[str], keywords: Sequence[str], fallback: str) -> str:
    if not sentences: return fallback
    ranked = sorted(sentences, key=lambda s: def_sentence_score(s, keywords, pattern), reverse=True)
    chosen = ranked[0] if def_sentence_score(ranked[0], keywords, pattern) > 0 else fallback
    return chosen[:DEFAULT_MAX_SENTENCE_CHARS].rstrip("。；; ")


def def_translate_fallback(label: str, zh: str) -> str:
    # Offline deterministic fallback: preserve facts; explicitly mark translation unavailable.
    return f"{label}: local English translation requires a GGUF or Transformers model; verified facts preserved in Chinese"


class DeterministicExtractiveSummarizer:
    name = "deterministic-extractive"
    def summarize(self, context: SummaryContext) -> list[str]:
        m = context.metrics
        valuation_parts = [x for x in [f"評等 {m.rating}" if m.rating != MISSING_VALUE else "",
                                       f"目標價 {m.target_price}" if m.target_price != MISSING_VALUE else "",
                                       f"上漲空間 {m.upside}" if m.upside != MISSING_VALUE else "",
                                       f"採 {m.valuation_method}" if m.valuation_method != MISSING_VALUE else ""] if x]
        valuation = "，".join(valuation_parts) if valuation_parts else "報告未提供可驗證的評等、目標價、上漲空間或評價方法"
        catalyst = def_pick_sentence(context.sentences, CATALYST_RE, ["需求", "市佔", "產能", "訂單", "AI", "HPC"], "報告未明確指出核心驅動")
        financial = def_pick_sentence(context.sentences, FINANCIAL_RE, ["營收", "毛利率", "EPS", "獲利", "成長"], "報告未提供可驗證的營運或財務亮點")
        risk = def_pick_sentence(context.sentences, RISK_RE, ["風險", "展望", "庫存", "匯率"], "報告未明確揭露風險；使用者應查閱原始報告與免責聲明")
        items = [valuation, catalyst, financial, risk]
        return [f"* 【{label}】{zh}。({def_translate_fallback(label, zh)}.)" for label, zh in zip(SUMMARY_LABELS, items)]


class LlamaCppSummarizer:
    name = "llama-cpp-python"
    def __init__(self, model_path: Path, threads: int, context_size: int):
        if not def_has_module("llama_cpp"): raise RuntimeError("請先安裝 llama-cpp-python")
        if not model_path.exists(): raise FileNotFoundError(model_path)
        from llama_cpp import Llama
        self.llm = Llama(model_path=str(model_path), n_ctx=context_size, n_threads=threads, verbose=False)

    def summarize(self, context: SummaryContext) -> list[str]:
        metrics_json = json.dumps(asdict(context.metrics), ensure_ascii=False)
        system = """你是金融研究報告摘要器。硬數據只可逐字使用 HARD_DATA，不得新增、推測或改寫數字。
僅輸出四行，每行以 * 開頭；順序與標籤必須是【評價與目標】【核心驅動】【營運與財務亮點】【風險與展望】。
每行格式必須為：* 【標籤】繁體中文摘要。(English translation.)
缺少資料必須寫「未提及」；不要開場、結語、Markdown 粗體或額外行。"""
        user = f"HARD_DATA={metrics_json}\nHEADER={context.header}\nREPORT_TEXT={context.clean_text[:DEFAULT_MAX_INPUT_CHARS]}"
        response = self.llm.create_chat_completion(messages=[{"role": "system", "content": system}, {"role": "user", "content": user}],
                                                   max_tokens=DEFAULT_MAX_TOKENS, temperature=DEFAULT_TEMPERATURE, top_p=DEFAULT_TOP_P)
        content = response["choices"][0]["message"]["content"].strip()
        return [line.strip() for line in content.splitlines() if line.strip().startswith("*")]


def def_choose_summarizer(backend: str, model_path: Optional[Path], threads: int, context_size: int) -> SummarizerBackend:
    if backend in {"llama", "auto"} and model_path:
        try: return LlamaCppSummarizer(model_path, threads, context_size)
        except Exception:
            if backend == "llama": raise
            logging.exception("GGUF backend unavailable; falling back to deterministic extractive summarizer")
    return DeterministicExtractiveSummarizer()


# =============================================================================
# def 08 VALIDATION / ANTI-HALLUCINATION
# =============================================================================

def def_numbers(text: str) -> set[str]:
    return {x.replace(",", "") for x in re.findall(r"(?<!\w)[-+]?\d+(?:,\d{3})*(?:\.\d+)?%?", text)}


def def_validate_summary(header: str, bullets: Sequence[str], metrics: InvestmentMetrics, source_text: str) -> tuple[str, list[str]]:
    findings: list[str] = []
    if len(bullets) != 4: findings.append(f"摘要必須恰為 4 點，目前為 {len(bullets)} 點")
    for idx, label in enumerate(SUMMARY_LABELS):
        if idx >= len(bullets) or f"【{label}】" not in bullets[idx]: findings.append(f"第 {idx + 1} 點標籤或順序錯誤：{label}")
        elif not re.search(r"\([^()]*[A-Za-z][^()]*\)\.?$", bullets[idx]): findings.append(f"第 {idx + 1} 點缺少括號英文")
    allowed = def_numbers(source_text + " " + json.dumps(asdict(metrics), ensure_ascii=False) + " 1 2 3 4")
    generated = def_numbers("\n".join(bullets))
    novel = sorted(generated - allowed)
    if novel: findings.append("摘要出現來源未驗證數字：" + ", ".join(novel))
    if not header or ":" not in header: findings.append("第一行標題格式錯誤")
    return ("PASS" if not findings else "REVIEW_REQUIRED"), findings


def def_repair_summary(context: SummaryContext, bullets: Sequence[str]) -> list[str]:
    gate, _ = def_validate_summary(context.header, bullets, context.metrics, context.clean_text)
    return list(bullets) if gate == "PASS" else DeterministicExtractiveSummarizer().summarize(context)


# =============================================================================
# def 09 ORCHESTRATION / OUTPUT
# =============================================================================

def def_blocks_to_clean_text(blocks: Sequence[DocumentBlock]) -> str:
    usable = [b.text_content for b in blocks if not b.excluded and b.sub_type not in {"頁首", "頁尾", "免責聲明", "小文字"}]
    return def_normalize_text("\n".join(usable))[:DEFAULT_MAX_INPUT_CHARS]


def def_run_engine(input_path: Path, security_master: Optional[Path] = None, ticker: str = "", backend: str = DEFAULT_BACKEND,
                   model_path: Optional[Path] = None, threads: int = DEFAULT_THREADS, context_size: int = DEFAULT_CONTEXT) -> EngineResult:
    if not input_path.exists(): raise FileNotFoundError(input_path)
    blocks, parser_name = def_parse_document(input_path)
    clean_text = def_repair_mojibake(def_blocks_to_clean_text(blocks))
    master = def_load_security_master(security_master)
    security = def_resolve_security(clean_text, master, ticker)
    raw_title = def_select_raw_title(blocks)
    title = def_clean_title(raw_title, security)
    header = def_build_header(title, security)
    metrics = def_extract_metrics(clean_text)
    sentences = def_dedupe_sentences(def_split_sentences(clean_text))
    context = SummaryContext(header, metrics, clean_text, sentences)
    summarizer = def_choose_summarizer(backend, model_path, threads, context_size)
    bullets = def_repair_summary(context, summarizer.summarize(context))
    gate, findings = def_validate_summary(header, bullets, metrics, clean_text)
    if security.status != "VERIFIED":
        findings.append("證券中文簡稱或市場未由官方／本機證券主檔驗證")
        gate = REVIEW_REQUIRED
    final_text = header + "\n" + "\n".join(bullets)
    stamp = datetime.now(timezone.utc)
    return EngineResult(stamp.strftime("RUN_%Y%m%dT%H%M%SZ"), str(input_path), def_sha256(input_path), ENGINE_NAME, ENGINE_VERSION,
                        stamp.isoformat(), parser_name, summarizer.name, security, header, metrics, bullets, final_text,
                        blocks, gate, findings, def_build_tool_registry())


def def_write_outputs(result: EngineResult, output_dir: Path) -> dict[str, str]:
    output_dir.mkdir(parents=True, exist_ok=True)
    stem = Path(result.source_file).stem
    paths = {"text": output_dir / f"{stem}_summary.txt", "json": output_dir / f"{stem}_result.json", "csv": output_dir / f"{stem}_blocks.csv"}
    paths["text"].write_text(result.final_text + "\n", encoding=DEFAULT_ENCODING)
    paths["json"].write_text(json.dumps(asdict(result), ensure_ascii=False, indent=2), encoding=DEFAULT_ENCODING)
    fields = list(asdict(result.blocks[0]).keys()) if result.blocks else list(DocumentBlock.__dataclass_fields__.keys())
    with paths["csv"].open("w", encoding=DEFAULT_ENCODING, newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fields); writer.writeheader()
        for block in result.blocks:
            row = asdict(block); row["metadata"] = json.dumps(row["metadata"], ensure_ascii=False); writer.writerow(row)
    if def_has_module("pyarrow") and result.blocks:
        import pyarrow as pa
        import pyarrow.parquet as pq
        parquet_path = output_dir / f"{stem}_blocks.parquet"
        pq.write_table(pa.Table.from_pylist([asdict(b) for b in result.blocks]), parquet_path)
        paths["parquet"] = parquet_path
    return {key: str(value) for key, value in paths.items()}


# =============================================================================
# def 10 SELF TEST / CLI
# =============================================================================

def def_self_test() -> None:
    sample = "台積電 (2330 TT) 先進製程供不應求，HPC需求續強。評等：買進，目標價850元，目前股價678元，上漲空間25.4%，採2027年25倍本益比。營收可望創高，毛利率54%。風險為匯率波動。"
    master = {"2330": def_make_security_identity("2330", "台積電", "TWSE", "TEST")}
    security = def_resolve_security(sample, master)
    assert security.yfinance_ticker == "2330.TW"
    title = def_clean_title(sample.split("。", 1)[0], security)
    assert "2330" not in title and "台積電" not in title
    metrics = def_extract_metrics(sample)
    assert metrics.target_price == "850元" and metrics.upside == "25.4%"
    ctx = SummaryContext(def_build_header(title, security), metrics, sample, def_split_sentences(sample))
    bullets = DeterministicExtractiveSummarizer().summarize(ctx)
    gate, findings = def_validate_summary(ctx.header, bullets, metrics, sample)
    assert gate == "PASS", findings

    # Regression matrix: market suffix must come from the master, never prefix heuristics.
    regression_master = {
        "2317": def_make_security_identity("2317", "鴻海", "TWSE", "TEST"),
        "5347": def_make_security_identity("5347", "世界先進", "TPEX", "TEST"),
        "0050": def_make_security_identity("0050", "元大台灣50", "TWSE", "TEST"),
    }
    header_cases = [
        ("【個股報告】鴻海(2317.TWO) 擴大 AI 伺服器佈局", "2317", "鴻海(2317.TW): 擴大 AI 伺服器佈局"),
        ("5347世界先進——成熟製程面臨價格壓力", "5347", "世界先進(5347.TWO): 成熟製程面臨價格壓力"),
        ("元大台灣50 (0050) 迎來除息行情", "0050", "元大台灣50(0050.TW): 迎來除息行情"),
    ]
    for raw, expected_ticker, expected_header in header_cases:
        identity = def_resolve_security(raw, regression_master)
        assert identity.ticker == expected_ticker
        actual = def_build_header(def_clean_title(raw, identity), identity)
        assert actual == expected_header, (actual, expected_header)

    # False-positive guards discovered by long TXT/DOCX/Markdown test corpora.
    assert def_resolve_security("2026 年報告，n_ctx=2048，純文件測試", {}).ticker == ""
    assert def_resolve_security("台積電 2330.TW 研究更新", {}).ticker == "2330"


def def_parse_args(argv: Optional[Sequence[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=ENGINE_NAME)
    parser.add_argument("input", nargs="?", type=Path, help="輸入文件")
    parser.add_argument("--output-dir", type=Path, default=Path("via_summarizer_output"))
    parser.add_argument("--security-master", type=Path, help="CSV/JSON；欄位 Ticker,Name,Market,Source")
    parser.add_argument("--ticker", default="", help="覆寫或補充股票代碼")
    parser.add_argument("--backend", choices=["auto", "extractive", "llama"], default=DEFAULT_BACKEND)
    parser.add_argument("--model-path", type=Path, help="本地 GGUF 路徑")
    parser.add_argument("--threads", type=int, default=DEFAULT_THREADS)
    parser.add_argument("--context-size", type=int, default=DEFAULT_CONTEXT)
    parser.add_argument("--self-test", action="store_true")
    return parser.parse_args(argv)


def def_main(argv: Optional[Sequence[str]] = None) -> int:
    args = def_parse_args(argv)
    logging.basicConfig(level=logging.INFO, format="def [%(levelname)s] %(message)s")
    if args.self_test:
        def_self_test(); print("def SELF_TEST: PASS")
        if not args.input: return 0
    if not args.input:
        print("def ERROR: 請提供輸入文件或使用 --self-test", file=sys.stderr); return 2
    try:
        result = def_run_engine(args.input, args.security_master, args.ticker, args.backend, args.model_path, args.threads, args.context_size)
        paths = def_write_outputs(result, args.output_dir)
        print(result.final_text)
        print(f"def QUALITY_GATE: {result.quality_gate}")
        print("def OUTPUTS: " + json.dumps(paths, ensure_ascii=False))
        return 0 if result.quality_gate == "PASS" else 3
    except Exception as exc:
        logging.exception("Engine failed: %s", exc); return 1


if __name__ == "__main__":
    raise SystemExit(def_main())
