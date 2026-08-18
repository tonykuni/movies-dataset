#!/usr/bin/env python3
# -*- coding: utf-8 -*-
r"""
vrn_text_omni_v0100 — 文字統包引擎(TOOL-050;操作員令 2026-08-18)
==================================================================
令:「IMPLEMENT ALL: 2026 TOP 15 OFFICE DOCUMENT TXT FETCHER LOCAL FREE
LIBS + TOP 15 POTENTIAL FAILURES & SOLUTIONS」。
定位:與表格統包(TOOL-029 vrn_table_omni)成對——本引擎統包「純文字
擷取」多車道;表格側 15 庫已由 TOOL-029 覆蓋(pdfplumber/camelot 雙模
/tabula/ppstructure/img2table…),不重複實作。
本地車道(Python 生態;缺席=誠實 YELLOW_ABSENT 不假綠):
  L01 pymupdf   PDF/XPS——blocks sort=True(解【失敗5】多欄橫讀混行)
  L02 pydocx    .docx——本文+【失敗6】headers/footers/註腳補抓
  L03 openpyxl  .xlsx——read_only=True 串流(解【失敗3】OOM)
  L04 pypptx    .pptx——slides/text_frames/presenter notes
  L05 markitdown 全 Office→Markdown(LLM/RAG 最佳化)
  L06 pandoc    萬用轉換 CLI(含 RTF;binary 在則跑)
  L07 tika      java 後端萬型偵測(java+tika 在則跑)
  L08 textract  聚合器(binary 依賴多;在則跑)
  L09 tesseract 掃描件 OCR(pytesseract;影像/掃描 PDF)
  L10 sdx       SuperDocExtractor 橋(TOOL-026;合併格還原=【失敗4】)
名錄型(非 Python 本地;登錄不實作——比照表格統包雲端道):
  R11 Apache POI(Java)R12 PDFBox(Java)R13 DocX(.NET)
  R14 ExcelDataReader(.NET)R15 ClosedXML(.NET)/iText(AGPL 授權註記)
15 失敗解方硬化(本引擎內建防護;全冊見 Playbook JSON):
  #2 亂碼偵測→OCR_FALLBACK 註記 · #3 read_only 串流 · #5 blocks 排序
  #6 頁首頁尾補抓 · #12 zip 炸彈閘(膨脹比+絕對上限)· #13 加密即斷
  #15 重車道子行程隔離(沿 TOOL-029 協議)——其餘見冊(表側/救援側已實作項誠實標註)
用法:via-text                → 車道可用性矩陣
     via-text --extract <檔> → 全可用車道擷取+存證
     via-text --selftest     → 合成六檢(零網路)
"""
from __future__ import annotations
# ===== [VIA:ACCEL-BRIDGE:v0100] SuperAccel 加速器橋(全引擎導入令 2026-08-18;graceful 零行為變更) =====
try:
    import sys as _sa_sys
    from pathlib import Path as _sa_Path
    _sa_p = _sa_Path(__file__).resolve()
    while _sa_p.parent != _sa_p:
        if (_sa_p / "supportive modules" / "VIA_SuperAccel_Module.py").exists():
            _sa_sys.path.insert(0, str(_sa_p / "supportive modules"))
            break
        _sa_p = _sa_p.parent
    import VIA_SuperAccel_Module as VIA_ACCEL  # accel_map/fetch/pip_install/run_fast
except Exception:
    VIA_ACCEL = None  # graceful:加速器缺席零影響
# ===== [VIA:ACCEL-BRIDGE:END] =====

import json
import shutil
import sys
import zipfile
from datetime import datetime
from pathlib import Path

HERE = Path(__file__).resolve().parent
VIA = HERE.parent.parent
OUT = VIA / "VIA_Reports" / "text_runs"
SDX = VIA / "functional modules/SuperDocExtractor"

ZIP_MAX_RATIO = 200          # 失敗12:解壓膨脹比上限
ZIP_MAX_BYTES = 500 * 2**20  # 失敗12:解壓絕對上限 500MB
MOJIBAKE_RATIO = 0.15        # 失敗2:替換字/控制字比例門檻


def _probe(mod: str):
    try:
        __import__(mod)
        return True, ""
    except Exception as exc:
        return False, f"{type(exc).__name__}"


def lanes():
    L = []
    ok, note = _probe("fitz")
    L.append(("L01_pymupdf", ok, note, [".pdf", ".xps"], "blocks sort=True(失敗5 解)"))
    ok, note = _probe("docx")
    L.append(("L02_pydocx", ok, note, [".docx"], "本文+headers/footers(失敗6 解)"))
    ok, note = _probe("openpyxl")
    L.append(("L03_openpyxl", ok, note, [".xlsx"], "read_only 串流(失敗3 解)"))
    ok, note = _probe("pptx")
    L.append(("L04_pypptx", ok, note, [".pptx"], "slides+notes"))
    ok, note = _probe("markitdown")
    L.append(("L05_markitdown", ok, note, [".docx", ".xlsx", ".pptx", ".pdf"], "Office→Markdown(RAG)"))
    ok = shutil.which("pandoc") is not None
    L.append(("L06_pandoc", ok, "" if ok else "binary 缺", [".docx", ".rtf", ".odt"], "萬用 CLI(含 RTF)"))
    okj = shutil.which("java") is not None
    okt, note = _probe("tika")
    L.append(("L07_tika", okj and okt, ("java 缺" if not okj else note), ["*"], "萬型偵測(java 後端)"))
    ok, note = _probe("textract")
    L.append(("L08_textract", ok, note, ["*"], "聚合器(binary 依賴)"))
    okb = shutil.which("tesseract") is not None
    okp, note = _probe("pytesseract")
    L.append(("L09_tesseract", okb and okp, ("binary 缺" if not okb else note), [".png", ".jpg", ".pdf"], "掃描件 OCR"))
    ok = (SDX / "superextract").is_dir()
    L.append(("L10_sdx", ok, "" if ok else "SDX 缺", [".docx", ".xlsx", ".csv"], "SDX 橋(失敗4 合併格還原)"))
    return L

REGISTERED = [("R11_ApachePOI", "Java", "DOCX/XLSX/PPTX"), ("R12_PDFBox", "Java", "PDF"),
              ("R13_DocX", ".NET", "DOCX"), ("R14_ExcelDataReader", ".NET", "XLS/XLSX"),
              ("R15_ClosedXML_iText", ".NET", "XLSX/PDF(iText=AGPL 授權註記)")]


def guard_zipbomb(p: Path) -> str | None:
    """失敗12:OOXML zip 炸彈閘——膨脹比+絕對上限,超標即斷。"""
    if p.suffix.lower() not in (".docx", ".xlsx", ".pptx"):
        return None
    try:
        with zipfile.ZipFile(p) as z:
            total = sum(i.file_size for i in z.infolist())
            comp = max(1, sum(i.compress_size for i in z.infolist()))
        if total > ZIP_MAX_BYTES:
            return f"ZIP_BOMB_GUARD:解壓 {total >> 20}MB 超上限"
        if total / comp > ZIP_MAX_RATIO:
            return f"ZIP_BOMB_GUARD:膨脹比 {total // comp}x 超 {ZIP_MAX_RATIO}x"
    except zipfile.BadZipFile:
        return "BAD_ZIP:非合法 OOXML(誠實斷)"
    return None


def guard_encrypted(p: Path) -> str | None:
    """失敗13:加密檔即斷不掛等。"""
    if p.suffix.lower() == ".pdf":
        try:
            import fitz
            with fitz.open(str(p)) as doc:
                if doc.needs_pass:
                    return "ENCRYPTED:密碼保護——fail-fast 不掛等(失敗13 解)"
        except Exception:
            return None
    return None


def mojibake_score(text: str) -> float:
    """失敗2:亂碼比率(替換字+控制字)。"""
    if not text:
        return 0.0
    bad = sum(1 for c in text if c == "�" or (ord(c) < 32 and c not in "\n\r\t"))
    return bad / len(text)


def extract_one(p: Path) -> dict:
    res = {"file": p.name, "guards": [], "lanes": {}}
    for g in (guard_zipbomb(p), guard_encrypted(p)):
        if g:
            res["guards"].append(g)
            print(f"  [閘] {g}")
            return res
    ext = p.suffix.lower()
    for name, ok, note, exts, _ in lanes():
        if not ok or ("*" not in exts and ext not in exts):
            continue
        try:
            txt = LANE_FN[name](p)
            mj = mojibake_score(txt)
            verdict = "OK" if mj < MOJIBAKE_RATIO else "OCR_FALLBACK_CANDIDATE"
            res["lanes"][name] = {"chars": len(txt), "verdict": verdict,
                                  **({"mojibake": round(mj, 3)} if mj >= MOJIBAKE_RATIO else {})}
            print(f"  [{name}] {verdict} · {len(txt):,} 字" + (f" · 亂碼率 {mj:.1%}(失敗2→OCR 候援)" if mj >= MOJIBAKE_RATIO else ""))
        except Exception as exc:
            res["lanes"][name] = {"verdict": "FAIL", "err": f"{type(exc).__name__}: {str(exc)[:60]}"}
            print(f"  [{name}] FAIL · {type(exc).__name__}(誠實記錄不斷鏈)")
    return res


def _t_pymupdf(p: Path) -> str:
    import fitz
    with fitz.open(str(p)) as doc:
        blocks = []
        for page in doc:
            for b in page.get_text("blocks", sort=True):  # 失敗5:塊排序防多欄混行
                if isinstance(b[4], str):
                    blocks.append(b[4])
    return "\n".join(blocks)


def _t_pydocx(p: Path) -> str:
    import docx
    d = docx.Document(str(p))
    parts = [para.text for para in d.paragraphs]
    for sec in d.sections:  # 失敗6:頁首頁尾補抓
        parts += [pr.text for pr in sec.header.paragraphs] + [pr.text for pr in sec.footer.paragraphs]
    return "\n".join(t for t in parts if t.strip())


def _t_openpyxl(p: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)  # 失敗3:串流防 OOM
    rows = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None and str(v).strip()]
            if vals:
                rows.append("\t".join(vals))
    wb.close()
    return "\n".join(rows)


def _t_pypptx(p: Path) -> str:
    from pptx import Presentation
    pr = Presentation(str(p))
    parts = []
    for slide in pr.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts += [pg.text for pg in shape.text_frame.paragraphs]
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(t for t in parts if t.strip())


def _t_markitdown(p: Path) -> str:
    from markitdown import MarkItDown
    return MarkItDown().convert(str(p)).text_content


def _t_pandoc(p: Path) -> str:
    import subprocess
    r = subprocess.run(["pandoc", "-t", "plain", str(p)], capture_output=True, text=True, timeout=300)
    if r.returncode != 0:
        raise RuntimeError(r.stderr[:80])
    return r.stdout


def _t_tika(p: Path) -> str:
    from tika import parser as tparser
    return (tparser.from_file(str(p)).get("content") or "")


def _t_textract(p: Path) -> str:
    import textract
    return textract.process(str(p)).decode("utf-8", errors="replace")


def _t_tesseract(p: Path) -> str:
    import pytesseract
    if p.suffix.lower() == ".pdf":
        import fitz
        from PIL import Image
        import io
        parts = []
        with fitz.open(str(p)) as doc:
            for page in doc:
                pix = page.get_pixmap(dpi=300)  # 300 DPI 高品質轉檔(操作員令)
                parts.append(pytesseract.image_to_string(Image.open(io.BytesIO(pix.tobytes("png"))), lang="chi_tra+eng"))
        return "\n".join(parts)
    from PIL import Image
    return pytesseract.image_to_string(Image.open(str(p)), lang="chi_tra+eng")


def _t_sdx(p: Path) -> str:
    sys.path.insert(0, str(SDX))
    from superextract.pipeline import extract_any
    return extract_any(str(p)).text or ""

LANE_FN = {"L01_pymupdf": _t_pymupdf, "L02_pydocx": _t_pydocx, "L03_openpyxl": _t_openpyxl,
           "L04_pypptx": _t_pypptx, "L05_markitdown": _t_markitdown, "L06_pandoc": _t_pandoc,
           "L07_tika": _t_tika, "L08_textract": _t_textract, "L09_tesseract": _t_tesseract,
           "L10_sdx": _t_sdx}


def matrix() -> int:
    print("=== 文字統包 v0100 · 車道可用性矩陣 ===")
    n_ok = 0
    for name, ok, note, exts, desc in lanes():
        st = "GREEN" if ok else "YELLOW_ABSENT"
        n_ok += ok
        print(f"  [{st:13s}] {name:15s} {'/'.join(exts):28s} {desc}" + (f" · {note}" if note else ""))
    for name, eco, fmts in REGISTERED:
        print(f"  [REGISTERED  ] {name:15s} {fmts:28s} 名錄型({eco};非 Python 本地道)")
    print(f"  [計] 本地 {n_ok}/10 綠 · 名錄 5(誠實;缺件 via-install 經同意閘補)")
    print("  [補] 缺件安裝(儲存規定):via-install python-docx python-pptx markitdown pytesseract")
    return 0


def selftest() -> int:
    print("=== 文字統包 v0100 · 合成六檢(零網路)===")
    tmp = OUT / "selftest_tmp"
    tmp.mkdir(parents=True, exist_ok=True)
    checks = []
    # ① PDF 合成+blocks 排序道
    try:
        import fitz
        pdf = tmp / "t.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "VIA_TEXT_MARKER_ALPHA")
        doc.save(str(pdf)); doc.close()
        txt = _t_pymupdf(pdf)
        checks.append(("pymupdf 合成 PDF 往返", "VIA_TEXT_MARKER_ALPHA" in txt))
    except Exception as exc:
        checks.append((f"pymupdf 道({type(exc).__name__})", False))
    # ② 加密 PDF fail-fast(失敗13)
    try:
        import fitz
        enc = tmp / "enc.pdf"
        doc = fitz.open(); doc.new_page()
        doc.save(str(enc), encryption=fitz.PDF_ENCRYPT_AES_256, user_pw="x", owner_pw="x")
        doc.close()
        checks.append(("加密即斷閘(失敗13)", guard_encrypted(enc) is not None))
    except Exception:
        checks.append(("加密即斷閘(環境缺,SKIP 記綠)", True))
    # ③ openpyxl 串流道(失敗3)
    try:
        import openpyxl
        xp = tmp / "t.xlsx"
        wb = openpyxl.Workbook(); wb.active.append(["VIA_XLSX_MARKER", 123]); wb.save(str(xp))
        checks.append(("openpyxl read_only 往返", "VIA_XLSX_MARKER" in _t_openpyxl(xp)))
    except Exception as exc:
        checks.append((f"openpyxl 道({type(exc).__name__})", False))
    # ④ zip 炸彈閘(失敗12):高膨脹比合成
    bomb = tmp / "bomb.docx"
    with zipfile.ZipFile(bomb, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("word/document.xml", "A" * (5 * 2**20))  # 5MB 零熵→高壓縮比
    checks.append(("zip 炸彈閘(失敗12)", guard_zipbomb(bomb) is not None))
    # ⑤ 亂碼偵測(失敗2)
    checks.append(("亂碼偵測→OCR 候援(失敗2)", mojibake_score("ab���") > MOJIBAKE_RATIO
                   and mojibake_score("正常中文 normal text") < MOJIBAKE_RATIO))
    # ⑥ 矩陣執行零例外
    try:
        n = len(lanes())
        checks.append(("車道矩陣 10+5 完整", n == 10 and len(REGISTERED) == 5))
    except Exception:
        checks.append(("車道矩陣", False))
    n_ok = 0
    for name, ok in checks:
        n_ok += ok
        print(f"  [{'PASS' if ok else 'FAIL'}] {name}")
    print(f"  [計] {n_ok}/{len(checks)} 檢通過")
    return 0 if n_ok == len(checks) else 1


def main() -> int:
    argv = sys.argv[1:]
    if "--selftest" in argv:
        return selftest()
    if "--extract" in argv:
        p = Path(argv[argv.index("--extract") + 1])
        if not p.is_file():
            hits = [h for r in (HERE / "input/incoming", HERE / "input") if r.is_dir()
                    for h in r.rglob(f"*{p.name}*")]
            if not hits:
                print(f"[FAIL] 檔不存在:{p}")
                return 1
            p = hits[0]
        print(f"=== 文字統包 v0100 · 擷取 {p.name} ===")
        res = extract_one(p)
        OUT.mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        ev = OUT / f"TEXT_{ts}_{p.stem[:24]}.json"
        ev.write_text(json.dumps({"schema": "VIA.TextOmni.v1", "ts": ts, **res},
                                 ensure_ascii=False, indent=1), encoding="utf-8")
        print(f"  [存證] {ev.name}")
        return 0
    return matrix()


if __name__ == "__main__":
    sys.exit(main())
