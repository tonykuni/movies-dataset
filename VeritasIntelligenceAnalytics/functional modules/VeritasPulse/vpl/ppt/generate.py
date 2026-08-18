"""VPL-OUT · Project PPT generator (python-pptx).
Turns a project dict (from core.store) + optimized charts into a polished
VIA Visual Lock deck. Target: full project presentation in seconds.
"""
# ===== [VIA:ACCEL-BRIDGE:v0100] SuperAccel 加速器橋(全引擎導入令 2026-08-18;graceful 零行為變更) =====
try:
    import sys as _sa_sys
    from pathlib import Path as _sa_Path
    _sa_p = _sa_Path(__file__).resolve()
    while _sa_p.parent != _sa_p:
        if (_sa_p / "supportive modules" / "VIA_SuperAccel_Module.py").exists():
            _sa_sys.path.insert(0, str(_sa_p / "supportive modules"))
            break
        _sa_p = _sa_p.parent
    import VIA_SuperAccel_Module as VIA_ACCEL  # accel_map/fetch/pip_install/run_fast
except Exception:
    VIA_ACCEL = None  # graceful:加速器缺席零影響
# ===== [VIA:ACCEL-BRIDGE:END] =====
import os
from datetime import datetime, date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from .. import theme as T
from ..charts import optimize as C

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _rgb(hexs):
    return RGBColor.from_string(hexs)


def _fill(shape, hexs):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hexs)
    shape.line.fill.background()


def _text(box, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    first = True
    for r in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = r.get("align", align)
        p.line_spacing = r.get("space", space)
        if "before" in r:
            p.space_before = Pt(r["before"])
        run = p.add_run()
        run.text = r["t"]
        f = run.font
        f.size = Pt(r.get("sz", 16))
        f.bold = r.get("b", False)
        f.italic = r.get("i", False)
        f.name = r.get("font", T.FONT_BODY)
        f.color.rgb = _rgb(r.get("c", T.INK))


def _dot(slide, x, y, size=0.14, color=T.ACCENT):
    d = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(size), Inches(size))
    _fill(d, color)
    d.adjustments[0] = 0.3
    return d


def _bg(slide, hexs):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _fill(r, hexs)
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------- slides ----------
def title_slide(prs, p):
    s = _blank(prs); _bg(s, T.INK)
    _dot(s, 0.9, 1.55, 0.34, T.ACCENT)
    _text(s.shapes.add_textbox(Inches(1.45), Inches(1.4), Inches(11), Inches(0.6)),
          [{"t": "VERITASPULSE", "sz": 14, "font": T.FONT_MONO, "c": T.INK_FAINT, "b": True}])
    _text(s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.4)),
          [{"t": p["name"], "sz": 46, "b": True, "font": T.FONT_DISPLAY, "c": T.BG, "space": 1.02}])
    _text(s.shapes.add_textbox(Inches(0.95), Inches(5.2), Inches(11), Inches(1.2)),
          [{"t": p["summary"], "sz": 14, "c": T.LINE, "space": 1.25}])
    _text(s.shapes.add_textbox(Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.5)),
          [{"t": f"{p['code']}   ·   Owner: {p['owner']}   ·   {p['start']} → {p['finish']}   ·   {date.today().isoformat()}",
            "sz": 11, "font": T.FONT_MONO, "c": T.INK_FAINT}])


def _kpi_card(s, x, y, w, label, value, sub, accent):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.9))
    _fill(card, T.PANEL); card.adjustments[0] = 0.06
    card.line.color.rgb = _rgb(T.LINE); card.line.width = Pt(1)
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + 0.18), Inches(0.07), Inches(1.54))
    _fill(strip, accent)
    _text(s.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.24), Inches(w - 0.4), Inches(0.4)),
          [{"t": label.upper(), "sz": 10, "font": T.FONT_MONO, "c": T.INK_FAINT}])
    _text(s.shapes.add_textbox(Inches(x + 0.26), Inches(y + 0.62), Inches(w - 0.4), Inches(0.9)),
          [{"t": value, "sz": 38, "b": True, "font": T.FONT_DISPLAY, "c": T.INK}])
    _text(s.shapes.add_textbox(Inches(x + 0.28), Inches(y + 1.46), Inches(w - 0.4), Inches(0.4)),
          [{"t": sub, "sz": 10.5, "font": T.FONT_MONO, "c": accent}])


def snapshot_slide(prs, p):
    s = _blank(prs); _bg(s, T.BG)
    _dot(s, 0.9, 0.78)
    _text(s.shapes.add_textbox(Inches(1.18), Inches(0.62), Inches(11), Inches(0.6)),
          [{"t": "Executive Snapshot", "sz": 26, "b": True, "font": T.FONT_DISPLAY, "c": T.INK}])
    tasks = p["tasks"]
    done = sum(1 for t in tasks if t["status"] == "done")
    pct = round(sum(t["progress"] for t in tasks) / len(tasks) * 100)
    planned = sum(b["planned"] for b in p["budget"])
    actual = sum(b["actual"] for b in p["budget"])
    open_risks = sum(1 for r in p["risks"] if r["status"] == "active")
    hit_ms = sum(1 for m in p["milestones"] if m["hit"])
    cards = [
        ("Completion", f"{pct}%", "weighted avg", T.ACCENT),
        ("Tasks Done", f"{done}/{len(tasks)}", "on track", T.BLUE),
        ("Budget Used", f"{actual/planned*100:.0f}%", f"${actual/1000:.0f}k / ${planned/1000:.0f}k", T.AMBER),
        ("Open Risks", f"{open_risks}", "active mitigations", T.RUST),
    ]
    x = 0.9; w = 2.86; gap = 0.18
    for lbl, val, sub, acc in cards:
        _kpi_card(s, x, 1.55, w, lbl, val, sub, acc); x += w + gap
    # status line
    _text(s.shapes.add_textbox(Inches(0.95), Inches(3.85), Inches(11.5), Inches(0.5)),
          [{"t": "MILESTONES", "sz": 11, "font": T.FONT_MONO, "c": T.INK_FAINT, "b": True}])
    y = 4.35
    for m in p["milestones"]:
        _dot(s, 1.0, y + 0.06, 0.12, T.ACCENT if m["hit"] else T.INK_FAINT)
        _text(s.shapes.add_textbox(Inches(1.3), Inches(y - 0.04), Inches(8), Inches(0.4)),
              [{"t": m["name"], "sz": 14, "b": True, "c": T.INK}])
        _text(s.shapes.add_textbox(Inches(9.4), Inches(y - 0.04), Inches(3), Inches(0.4)),
              [{"t": f"due {m['due']}   ·   {'hit' if m['hit'] else 'pending'}",
                "sz": 11, "font": T.FONT_MONO, "c": T.INK_FAINT, "align": PP_ALIGN.RIGHT}])
        y += 0.62


def chart_slide(prs, title, sub, img, *, img_w=8.0, side_runs=None):
    s = _blank(prs); _bg(s, T.BG)
    _dot(s, 0.9, 0.78)
    _text(s.shapes.add_textbox(Inches(1.18), Inches(0.62), Inches(11), Inches(0.6)),
          [{"t": title, "sz": 26, "b": True, "font": T.FONT_DISPLAY, "c": T.INK}])
    if sub:
        _text(s.shapes.add_textbox(Inches(1.18), Inches(1.18), Inches(11), Inches(0.4)),
              [{"t": sub, "sz": 12, "font": T.FONT_MONO, "c": T.INK_FAINT}])
    from PIL import Image
    iw, ih = Image.open(img).size
    avail_h = 5.45  # below the title block, above bottom margin
    full_width = side_runs is None
    avail_w = 11.53 if full_width else img_w
    disp_w = min(img_w, avail_w)
    disp_h = disp_w * ih / iw
    if disp_h > avail_h:                 # height-bound: rescale from height
        disp_h = avail_h
        disp_w = disp_h * iw / ih
    if full_width:                       # centre horizontally across the slide
        left = 0.9 + (avail_w - disp_w) / 2
    else:
        left = 0.9
    top = 1.7 + (avail_h - disp_h) / 2
    s.shapes.add_picture(img, Inches(left), Inches(top), width=Inches(disp_w))
    if side_runs:
        box = s.shapes.add_textbox(Inches(img_w + 1.1), Inches(1.9), Inches(13.333 - img_w - 2.0), Inches(5))
        _text(box, side_runs, space=1.15)
    return s


def closing_slide(prs, p):
    s = _blank(prs); _bg(s, T.INK)
    _dot(s, 0.9, 2.7, 0.34, T.ACCENT)
    _text(s.shapes.add_textbox(Inches(0.9), Inches(3.15), Inches(11.5), Inches(1.4)),
          [{"t": "Generated by VeritasPulse", "sz": 34, "b": True, "font": T.FONT_DISPLAY, "c": T.BG}])
    _text(s.shapes.add_textbox(Inches(0.95), Inches(4.5), Inches(11.5), Inches(0.6)),
          [{"t": f"VPL · Project & Time Intelligence   ·   {datetime.now():%Y-%m-%d %H:%M}",
            "sz": 12, "font": T.FONT_MONO, "c": T.INK_FAINT}])


def build(project, out_dir, out_name="VeritasPulse_ProjectDeck.pptx"):
    os.makedirs(out_dir, exist_ok=True)
    g = C.gantt(project["tasks"], out_dir)
    b = C.budget(project["budget"], out_dir)
    rk = C.risk_heatmap(project["risks"], out_dir)
    sk = C.stakeholders(project["stakeholders"], out_dir)

    prs = Presentation()
    prs.slide_width = EMU_W; prs.slide_height = EMU_H

    title_slide(prs, project)
    snapshot_slide(prs, project)
    chart_slide(prs, "Timeline & Schedule", "Gantt · status-coloured · progress overlay", g, img_w=11.4)

    planned = sum(x["planned"] for x in project["budget"])
    actual = sum(x["actual"] for x in project["budget"])
    var = planned - actual
    chart_slide(prs, "Budget", "Planned vs Actual", b, img_w=7.4, side_runs=[
        {"t": "VARIANCE", "sz": 11, "font": T.FONT_MONO, "c": T.INK_FAINT, "b": True},
        {"t": f"${var/1000:+.0f}k", "sz": 40, "b": True, "font": T.FONT_DISPLAY,
         "c": T.ACCENT if var >= 0 else T.RUST, "before": 4},
        {"t": f"{'Under' if var >= 0 else 'Over'} budget", "sz": 13, "c": T.INK_SOFT, "before": 2},
        {"t": f"Spent ${actual/1000:.0f}k of ${planned/1000:.0f}k", "sz": 12,
         "font": T.FONT_MONO, "c": T.INK_FAINT, "before": 14},
    ])

    top_risks = sorted(project["risks"], key=lambda r: r["probability"] * r["impact"], reverse=True)[:4]
    chart_slide(prs, "Risk & Crisis", "Probability × Impact", rk, img_w=6.0,
                side_runs=[{"t": "TOP RISKS", "sz": 11, "font": T.FONT_MONO, "c": T.INK_FAINT, "b": True}]
                + sum(([{"t": f"●  {r['name']}", "sz": 13, "b": True, "c": T.INK, "before": 10},
                        {"t": f"   P{r['probability']}×I{r['impact']} = {r['probability']*r['impact']}  ·  {r['status']}",
                         "sz": 11, "font": T.FONT_MONO, "c": T.INK_FAINT, "before": 1}] for r in top_risks), []))

    chart_slide(prs, "Stakeholders", "Power / Interest matrix", sk, img_w=6.4,
                side_runs=[{"t": "ENGAGEMENT", "sz": 11, "font": T.FONT_MONO, "c": T.INK_FAINT, "b": True}]
                + sum(([{"t": f"●  {s['name']}", "sz": 12.5, "b": True, "c": T.INK, "before": 9},
                        {"t": f"   {s['role']}  ·  P{s['power']}/I{s['interest']}",
                         "sz": 10.5, "font": T.FONT_MONO, "c": T.INK_FAINT, "before": 1}] for s in project["stakeholders"]), []))

    closing_slide(prs, project)

    out = os.path.join(out_dir, out_name)
    prs.save(out)
    return out
