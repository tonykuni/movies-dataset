"""附件擷取層：各類型文件 → 文字。

支援 pdf / docx / xlsx / pptx / txt / csv / md。
每個附件回傳 (text, meta)，meta 記錄擷取是否成功、字數、頁數 —— 供完整性驗證。
影像 OCR 為選用（環境有 pytesseract 才啟用），沒有則標記 skipped 不中斷。
"""
from __future__ import annotations

import os
from dataclasses import dataclass, field

from .encoding import detect_file_encoding, repair_mojibake


@dataclass
class AttachmentResult:
    path: str
    ext: str
    ok: bool
    text: str = ""
    chars: int = 0
    pages: int = 0
    reason: str = ""          # 失敗或略過原因（供完整性驗證）
    meta: dict = field(default_factory=dict)


def _pdf(path: str) -> AttachmentResult:
    import pdfplumber
    text_parts, npages = [], 0
    with pdfplumber.open(path) as pdf:
        npages = len(pdf.pages)
        for pg in pdf.pages:
            text_parts.append(pg.extract_text() or "")
            for tbl in (pg.extract_tables() or []):
                for row in tbl:
                    text_parts.append("\t".join(c or "" for c in row))
    txt = "\n".join(text_parts).strip()
    # 掃描型 PDF：沒文字 → 嘗試 OCR 救回（系統多做事）
    if not txt:
        ocr = _ocr_pdf(path)
        if ocr:
            return AttachmentResult(path, "pdf", ok=True, text=ocr, chars=len(ocr),
                                    pages=npages, meta={"recovered_by": "ocr"})
        return AttachmentResult(path, "pdf", ok=False, pages=npages,
                                reason="no_text_layer(掃描檔，需OCR)")
    return AttachmentResult(path, "pdf", ok=True, text=txt, chars=len(txt), pages=npages)


def _ocr_pdf(path: str) -> str:
    """掃描型 PDF 的 OCR 救援。需 pdf2image + pytesseract + tesseract；缺則回空字串。"""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        pages = convert_from_path(path, dpi=200)
        out = []
        for img in pages:
            out.append(pytesseract.image_to_string(img, lang="chi_tra+eng"))
        return "\n".join(out).strip()
    except Exception:
        return ""


def _rtf(path: str) -> AttachmentResult:
    enc = detect_file_encoding(path)
    with open(path, encoding=enc, errors="replace") as f:
        raw = f.read()
    # 極簡 RTF 去控制碼
    import re
    txt = re.sub(r"\\[a-z]+-?\d* ?", "", raw)
    txt = re.sub(r"[{}]", "", txt).strip()
    return AttachmentResult(path, "rtf", ok=bool(txt), text=txt, chars=len(txt))


def _html(path: str) -> AttachmentResult:
    enc = detect_file_encoding(path)
    with open(path, encoding=enc, errors="replace") as f:
        raw = f.read()
    import re
    txt = re.sub(r"<script.*?</script>", "", raw, flags=re.S | re.I)
    txt = re.sub(r"<style.*?</style>", "", txt, flags=re.S | re.I)
    txt = re.sub(r"<[^>]+>", " ", txt)
    txt = re.sub(r"\s+", " ", repair_mojibake(txt)).strip()
    return AttachmentResult(path, "html", ok=bool(txt), text=txt, chars=len(txt))


def _eml(path: str) -> AttachmentResult:
    import email
    from email import policy
    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f, policy=policy.default)
    subj = msg.get("subject", "")
    try:
        body = msg.get_body(preferencelist=("plain", "html"))
        content = body.get_content() if body else ""
    except Exception:
        content = ""
    txt = f"{subj}\n{repair_mojibake(content)}".strip()
    return AttachmentResult(path, "eml", ok=bool(txt), text=txt, chars=len(txt))


def _docx(path: str) -> AttachmentResult:
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for row in t.rows:
            parts.append("\t".join(c.text for c in row.cells))
    txt = "\n".join(x for x in parts if x).strip()
    return AttachmentResult(path, "docx", ok=bool(txt), text=txt, chars=len(txt),
                            reason="" if txt else "empty")


def _xlsx(path: str) -> AttachmentResult:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append("\t".join(cells))
    txt = "\n".join(parts).strip()
    return AttachmentResult(path, "xlsx", ok=bool(txt), text=txt, chars=len(txt))


def _pptx(path: str) -> AttachmentResult:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    txt = "\n".join(x for x in parts if x).strip()
    return AttachmentResult(path, "pptx", ok=bool(txt), text=txt, chars=len(txt),
                            pages=len(prs.slides.__iter__.__self__._sldIdLst))


def _text(path: str, ext: str) -> AttachmentResult:
    enc = detect_file_encoding(path)
    with open(path, encoding=enc, errors="replace") as f:
        txt = repair_mojibake(f.read()).strip()
    return AttachmentResult(path, ext, ok=bool(txt), text=txt, chars=len(txt))


def _image(path: str, ext: str) -> AttachmentResult:
    try:
        import pytesseract
        from PIL import Image
        txt = pytesseract.image_to_string(Image.open(path), lang="chi_tra+eng").strip()
        return AttachmentResult(path, ext, ok=bool(txt), text=txt, chars=len(txt),
                                reason="" if txt else "ocr_empty")
    except Exception as e:
        return AttachmentResult(path, ext, ok=False, reason=f"ocr_unavailable({e.__class__.__name__})")


_ROUTER = {
    ".pdf": _pdf, ".docx": _docx, ".xlsx": _xlsx, ".xls": _xlsx,
    ".pptx": _pptx, ".txt": lambda p: _text(p, "txt"),
    ".csv": lambda p: _text(p, "csv"), ".tsv": lambda p: _text(p, "tsv"),
    ".md": lambda p: _text(p, "md"), ".json": lambda p: _text(p, "json"),
    ".log": lambda p: _text(p, "log"), ".rtf": _rtf,
    ".html": _html, ".htm": _html, ".eml": _eml,
    ".png": lambda p: _image(p, "png"), ".jpg": lambda p: _image(p, "jpg"),
    ".jpeg": lambda p: _image(p, "jpeg"), ".tif": lambda p: _image(p, "tif"),
    ".tiff": lambda p: _image(p, "tiff"), ".bmp": lambda p: _image(p, "bmp"),
}


def extract_attachment(path: str) -> AttachmentResult:
    ext = os.path.splitext(path)[1].lower()
    fn = _ROUTER.get(ext)
    if fn is None:
        return AttachmentResult(path, ext.lstrip("."), ok=False, reason="unsupported_type")
    try:
        return fn(path)
    except Exception as e:
        return AttachmentResult(path, ext.lstrip("."), ok=False, reason=f"error:{e.__class__.__name__}:{e}")
