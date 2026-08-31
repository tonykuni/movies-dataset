"""產出層：簡報草稿（python-pptx）。

設計重點：
  - 有範本就沿用使用者範本（pptx_template）；沒有就用乾淨預設版面
  - 產出是『草稿』：分產品一頁 + 風險頁 + 總覽頁，供使用者修改定稿
  - 內容全部來自 SSOT，不杜撰

注意：這是草稿產生器，不是最終簡報。使用者看過、認可後才對外。
"""
from __future__ import annotations
# ===== [VIA:ACCEL-BRIDGE:v0100] SuperAccel 加速器橋(批102 全樹導入令;graceful 零行為變更) =====
try:
    import sys as _sa_sys
    from pathlib import Path as _sa_Path
    _sa_p = _sa_Path(__file__).resolve()
    while _sa_p.parent != _sa_p:
        if (_sa_p / "supportive modules" / "VIA_SuperAccel_Module.py").exists():
            _sa_sys.path.insert(0, str(_sa_p / "supportive modules"))
            break
        _sa_p = _sa_p.parent
    import VIA_SuperAccel_Module as VIA_ACCEL  # noqa: N816
except Exception:
    VIA_ACCEL = None  # graceful:加速器缺席零影響
# ===== [VIA:ACCEL-BRIDGE:END] =====

import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# 預設配色（沿用使用者既有 Visual Lock 慣例：紅=漲/警示、綠=完成）
_BG = RGBColor(0xF5, 0xF4, 0xF0)
_INK = RGBColor(0x2A, 0x2A, 0x2A)
_BLUE = RGBColor(0x4C, 0x78, 0xA8)
_RED = RGBColor(0xC9, 0x6B, 0x5A)
_GREEN = RGBColor(0x5A, 0x9E, 0x6F)
_GREY = RGBColor(0x9C, 0x98, 0x90)

_STATUS_COLOR = {"逾期": _RED, "已結": _GREEN, "進行中": _BLUE, "待辦": _GREY, "未標示": _GREY}


def _add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size, color, bold=False):
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Noto Sans CJK TC"


def build_deck(df: pd.DataFrame, stats: dict, out_path: str, template: str | None = None) -> str:
    prs = Presentation(template) if template else Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- 封面 ----
    s = _add_blank(prs)
    tf = _box(s, 0.8, 2.4, 11.7, 2)
    _set(tf.paragraphs[0], "本週工作進度", 40, _INK, bold=True)
    p = tf.add_paragraph(); _set(p, "SSOT 自動彙整 · 草稿（請複核後對外）", 18, _GREY)
    p = tf.add_paragraph()
    _set(p, f"新增 {stats['added']} · 變更 {stats['updated']} · 共 {stats['total']} 筆事件",
         14, _BLUE)

    # ---- 分產品頁 ----
    for prod, g in df.groupby("product"):
        prod_name = prod or "未分類"
        s = _add_blank(prs)
        head = _box(s, 0.7, 0.5, 12, 1)
        _set(head.paragraphs[0], f"{prod_name}", 28, _BLUE, bold=True)
        body = _box(s, 0.9, 1.7, 11.6, 5.3)
        first = True
        for _, r in g.iterrows():
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            st = r["status"] or "未標示"
            line = f"[{r['uid']}]  {r['title']}"
            _set(p, line, 16, _INK)
            tag = p.add_run()
            tag.text = f"   〔{st}〕"
            tag.font.size = Pt(13)
            tag.font.bold = True
            tag.font.color.rgb = _STATUS_COLOR.get(st, _GREY)
            if r["actor"] or r["topic"]:
                sub = body.add_paragraph()
                meta = " · ".join([x for x in [r["topic"], r["actor"]] if x])
                _set(sub, f"      {meta}", 12, _GREY)

    # ---- 風險頁 ----
    overdue = df[df["status"] == "逾期"]
    if len(overdue):
        s = _add_blank(prs)
        head = _box(s, 0.7, 0.5, 12, 1)
        _set(head.paragraphs[0], "⚠ 風險 / 需關注", 28, _RED, bold=True)
        body = _box(s, 0.9, 1.7, 11.6, 5.3)
        first = True
        for _, r in overdue.iterrows():
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            _set(p, f"[{r['uid']}]  {r['title']}  （{r['product'] or '未分類'}）", 16, _RED)

    prs.save(out_path)
    return out_path
