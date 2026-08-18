#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
================================================================================
VIA / VDF DocForge v2  ——  文件擷取・修復・分類引擎 (Fetch-Fix-Classify)
================================================================================
Veritas Intelligence Analytics · VeritasDataForge subsystem
Auto-ID: VIA-DOC-YYYYMMDD-######  |  Append-only ledger  |  UTF-8 no-BOM
SSOT registry: via_docforge_registry.json (只增不減)

v2 強化:
  · 全工具 lazy 掛載 + 啟動能力矩陣 (Celeritas 風格, 缺套件優雅降級)
  · 格式擴充: Word/Excel/PPT/Legacy/ODF/PDF/RTF/HTML/EML/MSG/Image/Google
  · OCR fallback (圖片 + 影像式 PDF, tesseract)
  · Excel 走 pandas 結構化摘要
  · Legacy .doc/.xls/.ppt 走 soffice headless 轉檔抽取
  · 內容去重 (BloomFilter size= 預篩 + sha256 精確確認)
  · SSOT 驅動分類 + PII 偵測 (TW 身分證/email → 自動升敏感度)
  · evidence 信心分層 (M/P/Est) + source_method 標記
  · @@PROGRESS|pct|label 進度串流 (stderr)
  · 內建 self-test harness:  python via_docforge.py --selftest

界線: 只讀「本機/已授權」磁碟檔; 不做雲端隱蔽抽取。stage-2 授權連接器見檔尾。
================================================================================
"""

import os
import re
import sys
import json
import zipfile
import hashlib
import shutil
import argparse
import subprocess
import tempfile
from datetime import datetime, timezone

UTC = timezone.utc
HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_REGISTRY = os.path.join(HERE, "via_docforge_registry.json")

# TextLab 文字智能層 (強化 reader / 斷句 / 錯別字修復 / 關鍵字 / 信件結構)
try:
    import via_docforge_textlab as _textlab
except Exception:
    _textlab = None


# ============================================================ 工具 lazy 註冊 ==
class ToolRegistry:
    """Celeritas 風格 lazy 載入: 首次使用才 import, 之後快取。缺套件回 None。"""
    _SPECS = {
        "docx":        ("docx",              "python-docx",        "Word 文字/表格抽取"),
        "openpyxl":    ("openpyxl",          "openpyxl",           "Excel 讀取"),
        "pandas":      ("pandas",            "pandas",             "Excel/CSV 結構化摘要"),
        "pptx":        ("pptx",              "python-pptx",        "PowerPoint 文字抽取"),
        "pikepdf":     ("pikepdf",           "pikepdf",            "PDF 修復/線性化"),
        "pypdf":       ("pypdf",             "pypdf",              "PDF 文字抽取"),
        "charset":     ("charset_normalizer","charset-normalizer", "編碼偵測"),
        "striprtf":    ("striprtf.striprtf", "striprtf",           "RTF 純文字"),
        "odf":         ("odf.opendocument",  "odfpy",              "OpenDocument 讀取"),
        "bs4":         ("bs4",               "beautifulsoup4",     "HTML 文字抽取"),
        "extract_msg": ("extract_msg",       "extract-msg",        "Outlook .msg 解析"),
        "olefile":     ("olefile",           "olefile",            "舊版 OLE 偵測"),
        "PIL":         ("PIL.Image",         "pillow",             "影像讀取"),
        "pytesseract": ("pytesseract",       "pytesseract",        "OCR (需 tesseract)"),
    }
    _BINS = {"tesseract": "OCR 引擎", "soffice": "LibreOffice headless 轉檔 (legacy office)"}

    def __init__(self):
        self._cache = {}

    def get(self, key):
        if key in self._cache:
            return self._cache[key]
        mod = None
        spec = self._SPECS.get(key)
        if spec:
            try:
                name = spec[0]
                mod = __import__(name, fromlist=["*"]) if "." in name else __import__(name)
            except Exception:
                mod = None
        self._cache[key] = mod
        return mod

    def has_bin(self, name):
        return shutil.which(name) is not None

    def matrix(self):
        m = {"libraries": {}, "binaries": {}}
        for k, (imp, pip, desc) in self._SPECS.items():
            m["libraries"][k] = {"available": self.get(k) is not None, "pip": pip, "desc": desc}
        for b, desc in self._BINS.items():
            m["binaries"][b] = {"available": self.has_bin(b), "desc": desc}
        return m


# =================================================== BloomFilter (size= 慣例) ==
class BloomFilter:
    """VDF v4.0 慣例: 建構參數用 size= 而非 capacity=。作為去重快速預篩。"""
    def __init__(self, size=1 << 20, hashes=4):
        self.size = int(size)
        self.k = int(hashes)
        self.bits = bytearray(self.size // 8 + 1)

    def _idx(self, item, i):
        h = hashlib.blake2b(f"{i}:{item}".encode("utf-8"), digest_size=8).digest()
        return int.from_bytes(h, "big") % self.size

    def add(self, item):
        for i in range(self.k):
            b = self._idx(item, i)
            self.bits[b // 8] |= (1 << (b % 8))

    def __contains__(self, item):
        for i in range(self.k):
            b = self._idx(item, i)
            if not (self.bits[b // 8] >> (b % 8) & 1):
                return False
        return True


class Dedup:
    """BloomFilter 預篩 + exact set 確認 (bloom 有 false-positive, 用 set 保正確)。"""
    def __init__(self):
        self.bloom = BloomFilter(size=1 << 20)
        self.exact = {}

    def check_add(self, content_hash, doc_id):
        if content_hash in self.bloom and content_hash in self.exact:
            return self.exact[content_hash]
        self.bloom.add(content_hash)
        self.exact[content_hash] = doc_id
        return None


# ================================================================== 常數 ==
FORMAT_MAP = {
    ".docx": "OFFICE_WORD", ".dotx": "OFFICE_WORD", ".docm": "OFFICE_WORD",
    ".xlsx": "OFFICE_EXCEL", ".xltx": "OFFICE_EXCEL", ".xlsm": "OFFICE_EXCEL",
    ".pptx": "OFFICE_PPT", ".potx": "OFFICE_PPT", ".pptm": "OFFICE_PPT",
    ".doc": "OFFICE_LEGACY", ".xls": "OFFICE_LEGACY", ".ppt": "OFFICE_LEGACY",
    ".odt": "ODF", ".ods": "ODF", ".odp": "ODF",
    ".pdf": "PDF", ".rtf": "RTF", ".html": "HTML", ".htm": "HTML",
    ".eml": "EMAIL_EML", ".msg": "EMAIL_MSG",
    ".png": "IMAGE", ".jpg": "IMAGE", ".jpeg": "IMAGE",
    ".tif": "IMAGE", ".tiff": "IMAGE", ".bmp": "IMAGE", ".gif": "IMAGE",
    ".gdoc": "GOOGLE_POINTER", ".gsheet": "GOOGLE_POINTER", ".gslides": "GOOGLE_POINTER",
    ".txt": "TEXT", ".csv": "TEXT", ".tsv": "TEXT",
    ".md": "TEXT", ".json": "TEXT", ".xml": "TEXT",
}
OOXML_FAMILIES = {"OFFICE_WORD", "OFFICE_EXCEL", "OFFICE_PPT"}
TEXTISH = {"TEXT", "HTML", "RTF", "GOOGLE_POINTER"}

PALETTE = {"bg": "#f5f4f0", "paper": "#ffffff", "ink": "#1e1d1a", "blue": "#4c78a8",
           "teal": "#439a9a", "up": "#c96b5a", "down": "#5a9e6f", "line": "#dbd9d3"}
HEALTH_COLOR = {"OK": PALETTE["down"], "REPAIRED": PALETTE["teal"],
                "DAMAGED": PALETTE["up"], "UNREADABLE": "#6b6a66", "DUPLICATE": PALETTE["blue"]}

RE_TWID = re.compile(r"\b[A-Z][12]\d{8}\b")
RE_EMAIL = re.compile(r"[\w.+-]+@[\w-]+\.[\w.-]+")
RE_PHONE = re.compile(r"\b09\d{2}[- ]?\d{3}[- ]?\d{3}\b")
FILENAME_BAD = ['\\', '/', ':', '*', '?', '"', '<', '>', '|', '\n', '\r', '\t']


# ================================================================ 工具函式 ==
def utcnow_iso():
    return datetime.now(UTC).strftime("%Y-%m-%dT%H:%M:%SZ")


def progress(pct, label):
    sys.stderr.write(f"@@PROGRESS|{int(pct)}|{label}\n")
    sys.stderr.flush()


def sanitize_filename(name, maxlen=80):
    for ch in FILENAME_BAD:
        name = name.replace(ch, "_")
    name = name.strip().strip(".")
    return (name[:maxlen] if name else "unnamed")


def gen_doc_id(seq, path):
    date = datetime.now(UTC).strftime("%Y%m%d")
    h = hashlib.blake2s(path.encode("utf-8"), digest_size=3).hexdigest().upper()
    return f"VIA-DOC-{date}-{seq:06d}", h


def sha256_of(path):
    try:
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(1 << 16), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return None


def norm_text_hash(text):
    t = re.sub(r"\s+", " ", (text or "")).strip().lower()
    return hashlib.sha256(t.encode("utf-8")).hexdigest() if t else None


def load_registry(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            reg = json.load(f)
        return reg["domains"], reg["sensitivity"]
    except Exception:
        return ({"FINANCE": ["invoice", "報價"], "TECHNICAL": ["spec", "架構"]},
                {"CONFIDENTIAL": ["機密", "confidential"], "INTERNAL": ["內部", "internal use"]})


def load_corrections(path):
    """從 SSOT registry 讀 append-only 修正字典 (context-free 對)。缺則回 None(用內建種子)。"""
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f).get("corrections") or None
    except Exception:
        return None


# ============================================================ 修復 (per family) ==
def repair_text_encoding(path, tools, out_dir):
    try:
        raw = open(path, "rb").read()
    except Exception as e:
        return "UNREADABLE", None, {"reason": f"read-fail: {e}"}
    detected, decoded = None, None
    cn = tools.get("charset")
    if cn:
        try:
            best = cn.from_bytes(raw).best()
            if best is not None:
                detected, decoded = best.encoding, str(best)
        except Exception:
            pass
    if decoded is None:
        for enc in ("utf-8-sig", "utf-8", "big5", "gbk", "cp950", "latin-1"):
            try:
                decoded, detected = raw.decode(enc), enc
                break
            except Exception:
                continue
    if decoded is None:
        decoded, detected = raw.decode("utf-8", errors="replace"), "utf-8/replace"
    nb = decoded.encode("utf-8")
    status = "REPAIRED" if nb != raw else "OK"
    fixed = None
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
        fixed = os.path.join(out_dir, sanitize_filename(os.path.basename(path)))
        open(fixed, "wb").write(nb)
    return status, fixed, {"detected_encoding": detected}


def repair_zip_container(path, out_dir, marker):
    """OOXML([Content_Types].xml) 與 ODF(mimetype) 通用 zip 完整性檢查/重打包。"""
    if not zipfile.is_zipfile(path):
        return "DAMAGED", None, {"reason": "not a zip container"}
    try:
        with zipfile.ZipFile(path) as z:
            bad = z.testzip()
            names = z.namelist()
    except Exception as e:
        return "DAMAGED", None, {"reason": f"zip open fail: {e}"}
    has_marker = marker in names
    if bad is None and has_marker:
        return "OK", None, {"members": len(names)}
    if not out_dir:
        return "DAMAGED", None, {"bad_member": bad, "has_marker": has_marker}
    os.makedirs(out_dir, exist_ok=True)
    fixed = os.path.join(out_dir, sanitize_filename(os.path.basename(path)))
    kept = dropped = 0
    try:
        with zipfile.ZipFile(path) as zin, zipfile.ZipFile(fixed, "w", zipfile.ZIP_DEFLATED) as zout:
            for n in zin.namelist():
                try:
                    zout.writestr(n, zin.read(n)); kept += 1
                except Exception:
                    dropped += 1
        status = "REPAIRED" if (kept and has_marker) else "DAMAGED"
        return status, (fixed if status == "REPAIRED" else None), {"kept": kept, "dropped": dropped}
    except Exception as e:
        return "DAMAGED", None, {"reason": f"repack fail: {e}"}


def repair_pdf(path, tools, out_dir):
    pk = tools.get("pikepdf")
    if pk:
        try:
            with pk.open(path) as pdf:
                pages, enc = len(pdf.pages), pdf.is_encrypted
            if out_dir:
                os.makedirs(out_dir, exist_ok=True)
                fixed = os.path.join(out_dir, sanitize_filename(os.path.basename(path)))
                with pk.open(path) as pdf:
                    pdf.save(fixed, linearize=True)
                return "REPAIRED", fixed, {"pages": pages, "encrypted": bool(enc)}
            return "OK", None, {"pages": pages, "encrypted": bool(enc)}
        except Exception as e:
            return "DAMAGED", None, {"reason": f"pikepdf: {e}"}
    pp = tools.get("pypdf")
    if pp:
        try:
            r = pp.PdfReader(path)
            return "OK", None, {"pages": len(r.pages)}
        except Exception as e:
            return "DAMAGED", None, {"reason": f"pypdf: {e}"}
    return "OK", None, {"note": "no pdf lib"}


def repair_image(path, tools, out_dir):
    if not tools.get("PIL"):
        return "OK", None, {"note": "no PIL"}
    try:
        from PIL import Image
        with Image.open(path) as im:
            im.verify()
        with Image.open(path) as im:
            return "OK", None, {"size": list(im.size), "mode": im.mode}
    except Exception as e:
        return "DAMAGED", None, {"reason": f"image: {e}"}


def repair_legacy(path, tools, out_dir):
    ole = tools.get("olefile")
    if ole:
        try:
            if not ole.isOleFile(path):
                return "DAMAGED", None, {"reason": "not a valid OLE (legacy) file"}
        except Exception as e:
            return "DAMAGED", None, {"reason": f"ole: {e}"}
    return "OK", None, {"note": "legacy OLE - passthrough"}


# ============================================================ 內容抽取 ==
def _ocr_image(path, tools):
    if not (tools.get("pytesseract") and tools.has_bin("tesseract") and tools.get("PIL")):
        return "", "NONE"
    try:
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(path), lang="chi_tra+eng").strip(), "OCR"
    except Exception:
        try:
            import pytesseract
            from PIL import Image
            return pytesseract.image_to_string(Image.open(path)).strip(), "OCR"
        except Exception:
            return "", "NONE"


def _ocr_pdf(path, tools):
    if not (tools.has_bin("tesseract") and shutil.which("pdftoppm")):
        return "", "NONE"
    try:
        import pytesseract
        from PIL import Image
        with tempfile.TemporaryDirectory() as td:
            subprocess.run(["pdftoppm", "-png", "-r", "150", "-l", "3", path,
                            os.path.join(td, "pg")], timeout=90,
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            buf = []
            for fn in sorted(os.listdir(td)):
                buf.append(pytesseract.image_to_string(Image.open(os.path.join(td, fn))))
            return "\n".join(buf).strip(), "OCR"
    except Exception:
        return "", "NONE"


def _soffice_to_txt(path, tools):
    if not tools.has_bin("soffice"):
        return "", "NONE"
    try:
        with tempfile.TemporaryDirectory() as td:
            subprocess.run(["soffice", "--headless", "--convert-to", "txt:Text",
                            "--outdir", td, path], timeout=60,
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            out = os.path.join(td, os.path.splitext(os.path.basename(path))[0] + ".txt")
            if os.path.exists(out):
                raw = open(out, "rb").read()
                for enc in ("utf-8", "big5", "gbk", "latin-1"):
                    try:
                        return raw.decode(enc), "SOFFICE"
                    except Exception:
                        continue
    except Exception:
        pass
    return "", "NONE"


def extract_content(path, family, tools, limit=8000):
    structured = {}
    # --- TextLab 強化 reader 優先 (Word/Excel/CSV/TXT) ---
    if _textlab:
        try:
            if family == "OFFICE_WORD" and tools.get("docx"):
                r = _textlab.Readers.read_word(path); return r["text"][:limit], "TEXTLAB", r["structured"]
            if family == "OFFICE_EXCEL" and (tools.get("pandas") or tools.get("openpyxl")):
                r = _textlab.Readers.read_excel(path); return r["text"][:limit], "TEXTLAB", r["structured"]
            if family == "TEXT" and path.lower().endswith((".csv", ".tsv")):
                r = _textlab.Readers.read_csv(path); return r["text"][:limit], "TEXTLAB", r["structured"]
            if family == "TEXT":
                r = _textlab.Readers.read_txt(path); return r["text"][:limit], "TEXTLAB", r["structured"]
        except Exception:
            pass  # 降級回原生抽取
    try:
        if family == "OFFICE_WORD" and tools.get("docx"):
            import docx
            d = docx.Document(path)
            parts = [p.text for p in d.paragraphs]
            for t in d.tables:
                for row in t.rows:
                    parts.append(" ".join(c.text for c in row.cells))
            return "\n".join(parts)[:limit], "NATIVE", structured

        if family == "OFFICE_EXCEL":
            pd = tools.get("pandas")
            if pd:
                xl = pd.ExcelFile(path)
                buf, sheets = [], []
                for sn in xl.sheet_names[:5]:
                    df = xl.parse(sn, nrows=40)
                    sheets.append({"name": sn, "rows": int(df.shape[0]), "cols": int(df.shape[1]),
                                   "columns": [str(c) for c in df.columns][:30]})
                    buf.append(sn + " " + " ".join(str(c) for c in df.columns))
                    buf.append(df.head(15).to_string(index=False))
                structured["sheets"] = sheets
                return "\n".join(buf)[:limit], "PANDAS", structured
            if tools.get("openpyxl"):
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                buf = []
                for ws in wb.worksheets[:3]:
                    for row in ws.iter_rows(max_row=40, values_only=True):
                        buf.append(" ".join("" if c is None else str(c) for c in row))
                return "\n".join(buf)[:limit], "NATIVE", structured

        if family == "OFFICE_PPT" and tools.get("pptx"):
            import pptx
            pr = pptx.Presentation(path)
            buf = []
            for s in pr.slides:
                for shp in s.shapes:
                    if shp.has_text_frame:
                        buf.append(shp.text_frame.text)
            return "\n".join(buf)[:limit], "NATIVE", structured

        if family == "OFFICE_LEGACY":
            txt, method = _soffice_to_txt(path, tools)
            return txt[:limit], method, structured

        if family == "ODF" and tools.get("odf"):
            from odf.opendocument import load
            from odf import text as odftext, teletype
            doc = load(path)
            buf = [teletype.extractText(p) for p in doc.getElementsByType(odftext.P)]
            return "\n".join(buf)[:limit], "NATIVE", structured

        if family == "PDF":
            txt = ""
            if tools.get("pypdf"):
                import pypdf
                r = pypdf.PdfReader(path)
                txt = "\n".join((pg.extract_text() or "") for pg in r.pages[:8])
            if len(txt.strip()) < 20:
                otxt, m = _ocr_pdf(path, tools)
                if otxt:
                    return otxt[:limit], m, structured
            return txt[:limit], "NATIVE", structured

        if family == "RTF" and tools.get("striprtf"):
            from striprtf.striprtf import rtf_to_text
            raw = open(path, "r", encoding="utf-8", errors="replace").read()
            return rtf_to_text(raw)[:limit], "NATIVE", structured

        if family == "HTML" and tools.get("bs4"):
            from bs4 import BeautifulSoup
            raw = open(path, "rb").read().decode("utf-8", "replace")
            parser = "lxml" if tools.get("bs4") and __import__("importlib").util.find_spec("lxml") else "html.parser"
            return BeautifulSoup(raw, parser).get_text(" ", strip=True)[:limit], "NATIVE", structured

        if family == "EMAIL_EML":
            import email
            from email import policy
            msg = email.message_from_binary_file(open(path, "rb"), policy=policy.default)
            subj = str(msg.get("subject", ""))
            frm = str(msg.get("from", ""))
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_content(); break
                if not body:
                    for part in msg.walk():
                        if part.get_content_type() == "text/html":
                            body = part.get_content(); break
            else:
                body = msg.get_content()
            structured["email"] = {"subject": subj, "from": frm}
            return f"{subj}\n{frm}\n{body}"[:limit], "EMAIL_PARSE", structured

        if family == "EMAIL_MSG" and tools.get("extract_msg"):
            import extract_msg
            m = extract_msg.Message(path)
            structured["email"] = {"subject": m.subject or "", "from": m.sender or ""}
            return f"{m.subject or ''}\n{m.sender or ''}\n{m.body or ''}"[:limit], "NATIVE", structured

        if family == "IMAGE":
            return (*_ocr_image(path, tools), structured)[0:2] + (structured,)

        if family in ("TEXT", "GOOGLE_POINTER"):
            raw = open(path, "rb").read()
            for enc in ("utf-8-sig", "utf-8", "big5", "gbk", "latin-1"):
                try:
                    return raw.decode(enc)[:limit], "NATIVE", structured
                except Exception:
                    continue
    except Exception as e:
        return "", f"ERR:{type(e).__name__}", structured
    return "", "NONE", structured


# ============================================================ 分類 ==
def detect_language(text):
    if not text:
        return "UNKNOWN"
    han = len(re.findall(r"[\u4e00-\u9fff]", text))
    latin = len(re.findall(r"[A-Za-z]", text))
    simp = len(re.findall(r"[这个说对国还与门问间际发]", text))
    trad = len(re.findall(r"[這個說對國還與門問間際發]", text))
    if han == 0 and latin > 0:
        return "EN"
    if han > 0 and latin > han:
        return "MIXED"
    if han > 0:
        return "ZH-HANS" if simp > trad else "ZH-HANT"
    return "UNKNOWN"


def _kw_score(text, name, table, default):
    hay = (name + "\n" + (text or "")).lower()
    scores = {lbl: sum(hay.count(k.lower()) for k in kws) for lbl, kws in table.items()}
    scores = {k: v for k, v in scores.items() if v}
    if not scores:
        return default, {}
    return max(scores, key=scores.get), scores


def classify(path, family, text, domains, sensitivity):
    name = os.path.basename(path)
    fmt_class = ("OFFICE" if family in OOXML_FAMILIES or family == "OFFICE_LEGACY"
                 else "GOOGLE" if family == "GOOGLE_POINTER"
                 else "EMAIL" if family in ("EMAIL_EML", "EMAIL_MSG")
                 else "IMAGE" if family == "IMAGE"
                 else family)
    domain, dscores = _kw_score(text, name, domains, "GENERAL")
    lang = detect_language(text)
    sens, sscores = _kw_score(text, name, sensitivity, "PUBLIC_OR_UNMARKED")
    pii = []
    if text:
        if RE_TWID.search(text):
            pii.append("TW_ID")
        if RE_EMAIL.search(text):
            pii.append("EMAIL")
        if RE_PHONE.search(text):
            pii.append("PHONE")
    if pii and sens == "PUBLIC_OR_UNMARKED":
        sens = "CONFIDENTIAL"
    conf = "M" if (text and (dscores or sscores)) else ("P" if text else "Est")
    return {"format_class": fmt_class, "content_domain": domain, "domain_scores": dscores,
            "language": lang, "sensitivity": sens, "pii": pii, "classify_confidence": conf}


# ============================================================ 單檔 pipeline ==
def process_file(path, out_root, seq, tools, dedup, domains, sensitivity, do_fix=True, corrections=None):
    ext = os.path.splitext(path)[1].lower()
    family = FORMAT_MAP.get(ext, "UNKNOWN")
    doc_id, phash = gen_doc_id(seq, path)
    try:
        st = os.stat(path); size = st.st_size
        mtime = datetime.fromtimestamp(st.st_mtime, UTC).strftime("%Y-%m-%dT%H:%M:%SZ")
    except Exception:
        size, mtime = None, None

    fixdir = os.path.join(out_root, "fixed", datetime.now(UTC).strftime("%Y%m%d")) if do_fix else ""

    if family in OOXML_FAMILIES:
        health, fixed, detail = repair_zip_container(path, fixdir, "[Content_Types].xml")
    elif family == "ODF":
        health, fixed, detail = repair_zip_container(path, fixdir, "mimetype")
    elif family == "PDF":
        health, fixed, detail = repair_pdf(path, tools, fixdir)
    elif family in TEXTISH:
        health, fixed, detail = repair_text_encoding(path, tools, fixdir)
    elif family == "IMAGE":
        health, fixed, detail = repair_image(path, tools, fixdir)
    elif family == "OFFICE_LEGACY":
        health, fixed, detail = repair_legacy(path, tools, fixdir)
    elif family in ("EMAIL_EML", "EMAIL_MSG"):
        health, fixed, detail = "OK", None, {"note": "email passthrough"}
    else:
        health, fixed, detail = "UNREADABLE", None, {"reason": "unsupported"}

    read_from = fixed if fixed else path
    text, method, structured = extract_content(read_from, family, tools) if family != "UNKNOWN" else ("", "NONE", {})
    # 分類跑在「修正後(繁體正規化)」文字上, 讓簡體/混寫內容能命中繁體關鍵字表
    classify_text = text
    if _textlab and text:
        try:
            classify_text = _textlab.normalize_and_correct(text, corrections=corrections)[0]
        except Exception:
            classify_text = text
    tags = classify(path, family, classify_text, domains, sensitivity)

    chash = norm_text_hash(text) or sha256_of(path)
    dup_of = dedup.check_add(chash, doc_id) if chash else None
    if dup_of:
        health = "DUPLICATE"

    # --- TextLab 富化: 斷句/錯別字修復/關鍵字 + 信件結構 ---
    text_intel, email_structure = None, None
    if _textlab and text:
        try:
            text_intel = _textlab.enrich_text(text, corrections=corrections)
            tags["keywords"] = [w for w, _ in text_intel.get("keywords", [])][:8]
        except Exception as e:
            text_intel = {"error": str(e)[:120]}
    if _textlab and family == "EMAIL_EML":
        try:
            email_structure = _textlab.analyze_email_file(path)
        except Exception as e:
            email_structure = {"error": str(e)[:120]}

    return {
        "doc_id": doc_id, "path_hash": phash,
        "source_path": os.path.abspath(path),
        "fixed_path": os.path.abspath(fixed) if fixed else None,
        "file_name": os.path.basename(path), "extension": ext,
        "format_family": family, "size_bytes": size, "modified_time": mtime,
        "sha256": sha256_of(path), "content_hash": chash,
        "health": health, "is_duplicate": bool(dup_of), "duplicate_of": dup_of,
        "repair_detail": detail, "source_method": method,
        "structured": structured, "aspects": tags,
        "text_intel": text_intel, "email_structure": email_structure,
        "processed_at": utcnow_iso(), "evidence": "LOCAL_AUTHORIZED_SCAN",
    }


# ============================================================ 帳本 / 輸出 ==
def load_ledger_count(p):
    if not os.path.exists(p):
        return 0
    return sum(1 for _ in open(p, "r", encoding="utf-8"))


def append_ledger(p, rec):
    os.makedirs(os.path.dirname(p), exist_ok=True)
    with open(p, "a", encoding="utf-8") as f:
        f.write(json.dumps(rec, ensure_ascii=False) + "\n")


def write_per_file_json(out_root, rec):
    d = os.path.join(out_root, "meta"); os.makedirs(d, exist_ok=True)
    fn = f"{rec['doc_id']}_{sanitize_filename(os.path.splitext(rec['file_name'])[0], 40)}.json"
    json.dump(rec, open(os.path.join(d, fn), "w", encoding="utf-8"), ensure_ascii=False, indent=2)


def build_html_report(records, caps, email_report=None):
    total = len(records)
    def cnt(fn):
        d = {}
        for r in records:
            k = fn(r); d[k] = d.get(k, 0) + 1
        return dict(sorted(d.items(), key=lambda kv: -kv[1]))
    by = {"健康 Health": cnt(lambda r: r["health"]),
          "格式 Format": cnt(lambda r: r["format_family"]),
          "領域 Domain": cnt(lambda r: r["aspects"]["content_domain"]),
          "語言 Language": cnt(lambda r: r["aspects"]["language"]),
          "敏感度 Sensitivity": cnt(lambda r: r["aspects"]["sensitivity"])}
    P = PALETTE
    def bars(title, data):
        mx = max(data.values()) if data else 1
        rows = "".join(f'<div class="row"><span class="lbl">{k}</span>'
                       f'<span class="bar" style="width:{int(v/mx*100)}%"></span>'
                       f'<span class="num">{v}</span></div>' for k, v in data.items())
        return f'<div class="card"><h3>{title}</h3>{rows}</div>'
    cap_rows = ""
    if caps:
        for k, info in list(caps["libraries"].items()) + list(caps["binaries"].items()):
            c = P["down"] if info["available"] else P["up"]
            cap_rows += f'<span class="chip" style="border-color:{c};color:{c}">{k}</span>'
    tr = ""
    for r in records[:500]:
        a = r["aspects"]; hc = HEALTH_COLOR.get(r["health"], "#888")
        pii = ",".join(a["pii"]) if a["pii"] else ""
        kws = "、".join(a.get("keywords", [])[:5])
        ti = r.get("text_intel") or {}
        ncorr = ti.get("correction_count", "")
        tr += (f'<tr><td class="mono">{r["doc_id"]}</td><td>{r["file_name"]}</td>'
               f'<td>{r["format_family"]}</td><td><b style="color:{hc}">{r["health"]}</b></td>'
               f'<td>{a["content_domain"]}</td><td>{a["language"]}</td>'
               f'<td>{a["sensitivity"]}</td><td style="color:{P["up"]}">{pii}</td>'
               f'<td>{kws}</td><td class="mono">{ncorr}</td>'
               f'<td class="mono">{r["source_method"]}</td><td class="mono">{a["classify_confidence"]}</td></tr>')
    dup = sum(1 for r in records if r["is_duplicate"])
    healthy = sum(1 for r in records if r["health"] in ("OK", "REPAIRED"))
    bad = sum(1 for r in records if r["health"] in ("DAMAGED", "UNREADABLE"))

    email_section = ""
    if email_report and email_report.get("threads"):
        th_rows = "".join(
            f'<tr><td>{t["subject"]}</td><td class="mono">{t["count"]}</td>'
            f'<td>{"、".join(t["participants"][:6])}</td></tr>'
            for t in email_report["threads"][:30])
        org_rows = "".join(
            f'<span class="chip" style="border-color:{P["blue"]};color:{P["blue"]}">'
            f'{d} ({len(people)})</span>' for d, people in
            sorted(email_report["organizations"].items(), key=lambda kv: -len(kv[1]))[:20])
        edge_rows = "".join(
            f'<tr><td>{e["from_org"]}</td><td>→</td><td>{e["to_org"]}</td>'
            f'<td class="mono">{e["weight"]}</td></tr>'
            for e in email_report["org_edges"][:20])
        email_section = f"""<h2 style="font-family:'Syne',sans-serif;margin:24px 0 8px">信件結構分析</h2>
<div class="card" style="margin-bottom:12px"><h3>組織 (依 email domain, 括號為人數)</h3>{org_rows}</div>
<div class="grid" style="grid-template-columns:1fr 1fr">
<div class="card"><h3>Thread 討論串 ({email_report['thread_count']})</h3>
<table><thead><tr><th>主題</th><th>封數</th><th>參與者</th></tr></thead><tbody>{th_rows}</tbody></table></div>
<div class="card"><h3>組織互動 (from → to, 權重)</h3>
<table><thead><tr><th>來源</th><th></th><th>目標</th><th>次數</th></tr></thead><tbody>{edge_rows}</tbody></table></div>
</div>"""
    return f"""<!doctype html><html lang="zh-Hant"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1"><title>VDF DocForge v2</title>
<style>*{{box-sizing:border-box}}body{{margin:0;background:{P['bg']};color:{P['ink']};
font-family:'DM Sans',system-ui,'Noto Sans TC',sans-serif;padding:28px}}
h1{{font-family:'Syne',sans-serif;margin:0 0 4px}}.sub{{color:#6b6a66;font-size:13px;margin-bottom:14px}}
.strip{{height:5px;border-radius:2px;margin:12px 0 20px;background:linear-gradient(90deg,{P['blue']},{P['teal']},{P['down']},#c9b45a,{P['up']},#8a6fb0,{P['blue']})}}
.grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:14px;margin-bottom:20px}}
.card{{background:{P['paper']};border:1px solid {P['line']};border-radius:3px;padding:14px 16px}}
.card h3{{margin:0 0 10px;font-size:13px;font-weight:600}}
.row{{display:flex;align-items:center;gap:8px;margin:5px 0;font-size:12px}}
.lbl{{width:120px;flex:0 0 120px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}}
.bar{{height:12px;background:{P['teal']};border-radius:2px;min-width:2px}}.num{{margin-left:auto;font-family:'DM Mono',monospace;color:#555}}
table{{width:100%;border-collapse:collapse;background:{P['paper']};border:1px solid {P['line']};font-size:12px}}
th,td{{padding:7px 9px;border-bottom:1px solid {P['line']};text-align:left}}th{{background:#eeece7;font-weight:600}}
.mono{{font-family:'DM Mono',monospace}}.kpi{{display:flex;gap:20px}}.kpi b{{font-size:26px;font-family:'Syne',sans-serif}}
.chip{{display:inline-block;border:1px solid;border-radius:10px;padding:1px 8px;margin:2px;font-size:11px;font-family:'DM Mono',monospace}}</style></head><body>
<h1>VDF DocForge v2 · 擷取修復分類儀表板</h1>
<div class="sub">VeritasIntelligenceAnalytics · LOCAL_AUTHORIZED_SCAN · {utcnow_iso()}</div>
<div class="strip"></div>
<div class="kpi card"><div>總文件<br><b>{total}</b></div>
<div>正常/修復<br><b style="color:{P['down']}">{healthy}</b></div>
<div>損毀/無法讀<br><b style="color:{P['up']}">{bad}</b></div>
<div>重複<br><b style="color:{P['blue']}">{dup}</b></div></div>
<div class="card" style="margin:14px 0"><h3>工具能力矩陣 Capability Matrix</h3>{cap_rows}</div>
<div class="grid">{''.join(bars(t, d) for t, d in by.items())}</div>
<table><thead><tr><th>DOC ID</th><th>檔名</th><th>格式</th><th>健康</th><th>領域</th>
<th>語言</th><th>敏感度</th><th>PII</th><th>關鍵字</th><th>修正</th><th>抽取</th><th>信心</th></tr></thead><tbody>{tr}</tbody></table>
{email_section}
<p class="sub" style="margin-top:12px">僅顯示前 500 筆; 完整見 meta/*.json 與 ledger.jsonl (append-only)。</p>
</body></html>"""


# ============================================================ 主流程 ==
def run(root, out_root, do_fix=True, limit=None, registry=DEFAULT_REGISTRY):
    tools = ToolRegistry()
    caps = tools.matrix()
    domains, sensitivity = load_registry(registry)
    corrections = load_corrections(registry)
    dedup = Dedup()
    ledger = os.path.join(out_root, "ledger.jsonl")
    seq = load_ledger_count(ledger)

    print(f"[DocForge v2] root={root}  out={out_root}  seq_start={seq}")
    avail = [k for k, v in caps["libraries"].items() if v["available"]]
    print(f"[DocForge v2] tools: {', '.join(avail)}")
    append_ledger(ledger, {"type": "RUN_HEADER", "started_at": utcnow_iso(),
                           "capability_matrix": caps, "registry": os.path.abspath(registry)})

    candidates = []
    for dp, _, files in os.walk(root):
        for fn in files:
            if os.path.splitext(fn)[1].lower() in FORMAT_MAP:
                candidates.append(os.path.join(dp, fn))
    candidates.sort()                      # deterministic 順序 → 連號可重現
    ncand = len(candidates)
    if limit:
        candidates = candidates[:limit]

    records = []
    for i, path in enumerate(candidates, 1):
        rec = process_file(path, out_root if do_fix else "", seq, tools, dedup,
                           domains, sensitivity, do_fix, corrections)
        append_ledger(ledger, rec)
        write_per_file_json(out_root, rec)
        records.append(rec)
        seq += 1
        progress(i / max(len(candidates), 1) * 100, f"{rec['doc_id']} {rec['health']}")
        print(f"  [{i}/{len(candidates)}] {rec['doc_id']} {rec['health']:<10} "
              f"{rec['format_family']:<13} {rec['aspects']['content_domain']:<10} {rec['file_name']}")

    # 信件結構: 彙整所有 email → thread 樹 + 組織圖
    email_report = None
    if _textlab:
        parses = [r["email_structure"] for r in records
                  if r.get("email_structure") and "error" not in r["email_structure"]]
        if parses:
            email_report = _textlab.build_threads(parses)
            json.dump(email_report, open(os.path.join(out_root, "email_structure.json"), "w",
                                         encoding="utf-8"), ensure_ascii=False, indent=2)
            print(f"[DocForge v2] 信件結構: {email_report['thread_count']} threads · "
                  f"{len(email_report['organizations'])} orgs → email_structure.json")

    report = os.path.join(out_root, "docforge_report.html")
    open(report, "w", encoding="utf-8").write(build_html_report(records, caps, email_report))
    print(f"\n[DocForge v2] 完成 {len(records)}/{ncand} 筆 · ledger={ledger} · report={report}")
    return records


# ============================================================ SELF-TEST ==
def run_selftest():
    tools = ToolRegistry()
    passed = failed = skipped = 0
    log = []
    def check(name, cond, skip=False):
        nonlocal passed, failed, skipped
        if skip:
            skipped += 1; log.append(f"  SKIP {name}"); return
        if cond:
            passed += 1; log.append(f"  PASS {name}")
        else:
            failed += 1; log.append(f"  FAIL {name}")

    work = tempfile.mkdtemp(prefix="docforge_test_")
    src = os.path.join(work, "src"); out = os.path.join(work, "out")
    os.makedirs(src, exist_ok=True)

    open(os.path.join(src, "報價單.txt"), "wb").write("親愛的客戶 報價單 invoice 營收".encode("big5"))
    open(os.path.join(src, "broken.docx"), "wb").write(b"NOT A ZIP")
    if tools.get("docx"):
        import docx
        d = docx.Document(); d.add_paragraph("架構 spec 部署 pipeline 機密")
        d.save(os.path.join(src, "tech.docx"))
    if tools.get("pandas"):
        import pandas as pd
        pd.DataFrame({"客戶": ["A", "B"], "訂單金額": [100, 200]}).to_excel(
            os.path.join(src, "sales.xlsx"), index=False)
    if tools.get("pptx"):
        import pptx
        pr = pptx.Presentation(); sl = pr.slides.add_slide(pr.slide_layouts[5])
        sl.shapes.title.text = "合約 contract 法遵"
        pr.save(os.path.join(src, "legal.pptx"))
    if tools.get("pikepdf"):
        import pikepdf
        p = pikepdf.new(); p.add_blank_page(page_size=(200, 200))
        p.save(os.path.join(src, "blank.pdf"))
    eml = ("Subject: 採購 交期 confidential\nFrom: buyer@ext.com\n"
           "Content-Type: text/plain; charset=utf-8\n\n客戶下單 身分證 A123456789\n")
    open(os.path.join(src, "mail.eml"), "wb").write(eml.encode("utf-8"))
    if tools.get("striprtf"):
        open(os.path.join(src, "note.rtf"), "w", encoding="utf-8").write(
            r"{\rtf1\ansi 薪資 payroll 人事 report}")
    if tools.get("bs4"):
        open(os.path.join(src, "page.html"), "w", encoding="utf-8").write(
            "<html><body><h1>營收 revenue</h1><p>財報 balance sheet</p></body></html>")
    if tools.get("odf"):
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        od = OpenDocumentText(); od.text.addElement(P(text="測試報告 spec 架構"))
        od.save(os.path.join(src, "doc.odt"))
    ocr_ok = bool(tools.get("PIL") and tools.get("pytesseract") and tools.has_bin("tesseract"))
    if tools.get("PIL"):
        from PIL import Image, ImageDraw
        im = Image.new("RGB", (360, 90), "white")
        ImageDraw.Draw(im).text((10, 35), "INVOICE 2026", fill="black")
        im.save(os.path.join(src, "scan.png"))
    if tools.get("docx"):
        import docx
        d = docx.Document(); d.add_paragraph("架構 spec 部署 pipeline 機密")
        d.save(os.path.join(src, "tech_copy.docx"))
    legacy_made = False
    if tools.has_bin("soffice") and tools.get("docx"):
        try:
            import docx
            d = docx.Document(); d.add_paragraph("legacy 合約 contract 授權書")
            tmpx = os.path.join(work, "leg.docx"); d.save(tmpx)
            subprocess.run(["soffice", "--headless", "--convert-to", "doc",
                            "--outdir", src, tmpx], timeout=90,
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            legacy_made = os.path.exists(os.path.join(src, "leg.doc"))
        except Exception:
            legacy_made = False

    recs = run(src, out, do_fix=True, registry=DEFAULT_REGISTRY)
    idx = {r["file_name"]: r for r in recs}
    A = lambda fn, k: idx.get(fn, {}).get("aspects", {}).get(k)

    check("Big5 文字 REPAIRED",        idx.get("報價單.txt", {}).get("health") == "REPAIRED")
    check("Big5 文字 FINANCE 領域",     A("報價單.txt", "content_domain") == "FINANCE")
    check("Big5 文字 ZH-HANT",         A("報價單.txt", "language") == "ZH-HANT")
    check("壞 docx DAMAGED",           idx.get("broken.docx", {}).get("health") == "DAMAGED")
    check("docx TECHNICAL 領域",       A("tech.docx", "content_domain") == "TECHNICAL", skip="tech.docx" not in idx)
    check("docx 敏感度 CONFIDENTIAL",   A("tech.docx", "sensitivity") == "CONFIDENTIAL", skip="tech.docx" not in idx)
    check("xlsx pandas 結構化",         bool(idx.get("sales.xlsx", {}).get("structured", {}).get("sheets")), skip="sales.xlsx" not in idx)
    check("xlsx SALES 領域",            A("sales.xlsx", "content_domain") == "SALES", skip="sales.xlsx" not in idx)
    check("pptx LEGAL 領域",            A("legal.pptx", "content_domain") == "LEGAL", skip="legal.pptx" not in idx)
    check("pdf 健康 OK/REPAIRED",       idx.get("blank.pdf", {}).get("health") in ("OK", "REPAIRED"), skip="blank.pdf" not in idx)
    check("eml 解析 SALES 領域",        A("mail.eml", "content_domain") == "SALES")
    check("eml PII(TW_ID)→CONFIDENTIAL","TW_ID" in (A("mail.eml", "pii") or []) and A("mail.eml", "sensitivity") == "CONFIDENTIAL")
    check("rtf HR 領域",                A("note.rtf", "content_domain") == "HR", skip="note.rtf" not in idx)
    check("html FINANCE 領域",          A("page.html", "content_domain") == "FINANCE", skip="page.html" not in idx)
    check("odt TECHNICAL 領域",         A("doc.odt", "content_domain") == "TECHNICAL", skip="doc.odt" not in idx)
    check("image OCR 抽取",             idx.get("scan.png", {}).get("source_method") == "OCR", skip=not ocr_ok)
    check("重複偵測 DUPLICATE",         idx.get("tech_copy.docx", {}).get("health") == "DUPLICATE", skip="tech.docx" not in idx)
    check("legacy .doc soffice 抽取",   idx.get("leg.doc", {}).get("source_method") == "SOFFICE", skip=not legacy_made)
    check("auto-ID 格式正確",           all(re.match(r"VIA-DOC-\d{8}-\d{6}", r["doc_id"]) for r in recs))
    check("append-only 帳本存在",       os.path.exists(os.path.join(out, "ledger.jsonl")))
    check("HTML 報表生成",              os.path.exists(os.path.join(out, "docforge_report.html")))
    check("能力矩陣建立",               "libraries" in tools.matrix())
    # --- TextLab 整合 ---
    tl = _textlab is not None
    check("text_intel 關鍵字擷取",      bool((idx.get("tech.docx", {}).get("text_intel") or {}).get("keywords")), skip=not tl or "tech.docx" not in idx)
    check("text_intel 斷句計數",        (idx.get("page.html", {}).get("text_intel") or {}).get("sentence_count", 0) >= 1, skip=not tl or "page.html" not in idx)
    check("email 結構解析(from)",       bool((idx.get("mail.eml", {}).get("email_structure") or {}).get("participants", {}).get("from")), skip=not tl)
    check("email_structure.json 產出",  os.path.exists(os.path.join(out, "email_structure.json")), skip=not tl)

    print("\n" + "=" * 58)
    print("\n".join(log))
    print("=" * 58)
    print(f"SELFTEST: {passed}/{passed+failed} PASS"
          + (f"  ({skipped} skipped, deps missing)" if skipped else "  (0 skipped)"))
    print("=" * 58)
    shutil.rmtree(work, ignore_errors=True)
    return failed == 0


def main():
    ap = argparse.ArgumentParser(description="VIA VDF DocForge v2")
    ap.add_argument("root", nargs="?", help="要掃描的本機/已授權根目錄")
    ap.add_argument("-o", "--out", default=r".\DocForgeOut")
    ap.add_argument("-r", "--registry", default=DEFAULT_REGISTRY)
    ap.add_argument("--no-fix", action="store_true")
    ap.add_argument("--limit", type=int, default=None)
    ap.add_argument("--caps", action="store_true", help="只印工具能力矩陣")
    ap.add_argument("--selftest", action="store_true", help="跑自我測試")
    args = ap.parse_args()

    if args.selftest:
        sys.exit(0 if run_selftest() else 1)
    if args.caps:
        print(json.dumps(ToolRegistry().matrix(), ensure_ascii=False, indent=2)); return
    if not args.root:
        ap.error("需要 root 目錄 (或用 --selftest / --caps)")
    os.makedirs(args.out, exist_ok=True)
    run(args.root, args.out, do_fix=not args.no_fix, limit=args.limit, registry=args.registry)


if __name__ == "__main__":
    main()

# ==============================================================================
# STAGE-2 (雲端來源檔連接器 — 授權正路, 本檔未實作):
#   M365 Graph (Files.Read/Mail.Read) 或 Google Drive API, 走白名單報備 + delta 增量,
#   附 business justification 留稽核軌跡; 落地後直接餵 run(),
#   evidence 改標 AUTHORIZED_API_PULL。本引擎不含任何節流偽裝/監理規避邏輯。
# ==============================================================================
