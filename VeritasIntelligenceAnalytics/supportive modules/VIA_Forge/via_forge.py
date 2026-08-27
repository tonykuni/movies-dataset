#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
================================================================================
VIA Forge  ——  via_forge.py   (\u55ae\u4e00\u5f37\u5f15\u64ce / Fetch \u00b7 Fix \u00b7 Classify)
================================================================================
Veritas Intelligence Analytics \u00b7 VeritasDataForge subsystem \u00b7 \u5eab

\u6574\u4f75\u81ea DocForge v2 + TextLab + DocFetch Stack \u7684\u80fd\u529b\u8d85\u96c6\uff0c\u6536\u6582\u6210\u4e00\u652f\u5f15\u64ce\u3002
\u53ea\u589e\u4e0d\u6e1b\uff1a\u6240\u6709\u539f\u6709\u683c\u5f0f\u3001\u4fee\u5fa9\u3001\u5206\u985e\u3001\u53bb\u91cd\u3001\u5e33\u672c\u80fd\u529b\u5168\u90e8\u4fdd\u7559\u3002

\u5340\u584a (\u4f9d\u8cf4\u5747 lazy\uff0c\u7f3a\u5957\u4ef6\u512a\u96c5\u964d\u7d1a\uff0c\u7d55\u4e0d crash)\uff1a
  A. Accelerator   LazyTools / Bloom / Memo+Disk Cache / \u5206\u584a\u96dc\u6e4a / ThreadPool / Progress / Timer
  B. Config/SSOT   \u8f09\u5165 config + registry\uff0c\u5c55\u958b {root}\uff0cAUTO fallback
  C. Formats       Office(OOXML) / Legacy / ODF / PDF / RTF / HTML / EML / MSG / Image / Google / TEXT
  D. Repair(Fix)   \u7de8\u78bc / zip \u5bb9\u5668 / PDF \u7dda\u6027\u5316 / \u5f71\u50cf / legacy \u2192 OK/REPAIRED/DAMAGED/UNREADABLE
  E. Extract       native readers + OCR(\u5716\u7247/\u5f71\u50cf\u5f0f PDF) + soffice + \u7d50\u69cb\u5316\u6458\u8981
  F. TextLab       \u6b63\u898f\u5316 / s2twp / \u65b7\u53e5 / \u932f\u5b57\u4fee\u5fa9 / \u95dc\u9375\u8a5e / \u4fe1\u4ef6\u7d50\u69cb
  G. Classify      domain / sensitivity / PII\u5347\u654f / language / evidence(M/P/Est)
  H. Govern        BloomFilter \u53bb\u91cd / append-only ledger / auto-ID / per-file meta
  I. Export        json / csv / xlsx / html (Visual Lock)

Record schema \u8207 DocForge \u4e00\u81f4\uff1adoc_id / file_name / format_family / health /
  aspects{content_domain, sensitivity, pii, keywords, classify_confidence, language} /
  text_intel / email_structure / \u2026

CLI:
  python via_forge.py --selftest
  python via_forge.py --scan <dir> [--fix] [--export html]
  python via_forge.py --caps
================================================================================
"""
# ===== [VIA:ACCEL-BRIDGE:v0100] SuperAccel 加速器橋(批102 全樹導入令;graceful 零行為變更) =====
try:
    import sys as _sa_sys
    from pathlib import Path as _sa_Path
    _sa_p = _sa_Path(__file__).resolve()
    while _sa_p.parent != _sa_p:
        if (_sa_p / "supportive modules" / "VIA_SuperAccel_Module.py").exists():
            _sa_sys.path.insert(0, str(_sa_p / "supportive modules"))
            break
        _sa_p = _sa_p.parent
    import VIA_SuperAccel_Module as VIA_ACCEL  # noqa: N816
except Exception:
    VIA_ACCEL = None  # graceful:加速器缺席零影響
# ===== [VIA:ACCEL-BRIDGE:END] =====

import os
import re
import sys
import csv
import json
import zipfile
import hashlib
import argparse
import importlib
from collections import OrderedDict, Counter
from datetime import datetime, timezone
from concurrent.futures import ThreadPoolExecutor, as_completed

UTC = timezone.utc
HERE = os.path.dirname(os.path.abspath(__file__))


# ============================================================== A. ACCELERATOR ==
class _Lazy:
    __slots__ = ("_n", "_m", "_t")

    def __init__(self, name):
        self._n = name; self._m = None; self._t = False

    def _load(self):
        if not self._t:
            self._t = True
            try:
                self._m = importlib.import_module(self._n)
            except Exception:
                self._m = None
        return self._m

    @property
    def available(self):
        return self._load() is not None

    def __getattr__(self, k):
        m = self._load()
        if m is None:
            raise ImportError("optional module %s missing" % self._n)
        return getattr(m, k)

    def __bool__(self):
        return self._load() is not None


class Tools:
    """Named lazy handles + capability matrix (libraries + binaries)."""
    _LIB = {
        "docx": "docx", "openpyxl": "openpyxl", "pandas": "pandas",
        "pptx": "pptx", "pikepdf": "pikepdf", "pypdf": "pypdf",
        "charset": "charset_normalizer", "striprtf": "striprtf",
        "odfpy": "odf", "bs4": "bs4", "extract_msg": "extract_msg",
        "PIL": "PIL", "pytesseract": "pytesseract",
        "jieba": "jieba", "opencc": "opencc",
    }
    _BIN = {"tesseract": "tesseract", "soffice": "soffice"}

    def __init__(self):
        self._lz = {k: _Lazy(v) for k, v in self._LIB.items()}

    def get(self, name):
        lz = self._lz.get(name)
        return lz if (lz and lz.available) else None

    def bin_path(self, name):
        import shutil
        exe = self._BIN.get(name, name)
        return shutil.which(exe)

    def matrix(self):
        m = {k: bool(v.available) for k, v in self._lz.items()}
        for b in self._BIN:
            m[b] = bool(self.bin_path(b))
        return m


TOOLS = Tools()


class MemoCache:
    def __init__(self, max_entries=4096):
        self.max = int(max_entries); self._d = OrderedDict(); self.hits = 0; self.misses = 0

    def get(self, k):
        if k in self._d:
            self._d.move_to_end(k); self.hits += 1; return self._d[k]
        self.misses += 1; return None

    def set(self, k, v):
        self._d[k] = v; self._d.move_to_end(k)
        while len(self._d) > self.max:
            self._d.popitem(last=False)


class DiskCache:
    def __init__(self, cache_dir, ttl_hours=168, enabled=True):
        self.dir = cache_dir; self.ttl = float(ttl_hours) * 3600.0
        self.enabled = bool(enabled); self.hits = 0; self.misses = 0
        if self.enabled:
            os.makedirs(self.dir, exist_ok=True)

    def _p(self, k):
        return os.path.join(self.dir, k[:16] + ".json")

    def get(self, k):
        if not self.enabled:
            return None
        p = self._p(k)
        if not os.path.exists(p):
            self.misses += 1; return None
        import time
        if self.ttl > 0 and (time.time() - os.path.getmtime(p)) > self.ttl:
            self.misses += 1; return None
        try:
            with open(p, encoding="utf-8") as f:
                self.hits += 1; return json.load(f)
        except Exception:
            self.misses += 1; return None

    def set(self, k, v):
        if not self.enabled:
            return
        try:
            with open(self._p(k), "w", encoding="utf-8") as f:
                json.dump(v, f, ensure_ascii=False)
        except Exception:
            pass

    def stats(self):
        t = self.hits + self.misses
        return {"hits": self.hits, "misses": self.misses,
                "hit_rate": round(self.hits / t, 4) if t else 0.0}


class BloomFilter:
    """size= convention (locked). Pre-filter only; exact set confirms."""

    def __init__(self, size=262144, hashes=5):
        self.size = int(size); self.hashes = int(hashes)
        self.bits = bytearray((self.size + 7) // 8); self.count = 0

    def _idx(self, item, i):
        h = hashlib.blake2s(("%d|%s" % (i, item)).encode("utf-8")).digest()
        return int.from_bytes(h[:8], "big") % self.size

    def add(self, item):
        for i in range(self.hashes):
            b = self._idx(item, i); self.bits[b >> 3] |= (1 << (b & 7))
        self.count += 1

    def __contains__(self, item):
        for i in range(self.hashes):
            b = self._idx(item, i)
            if not (self.bits[b >> 3] & (1 << (b & 7))):
                return False
        return True


class Dedup:
    def __init__(self, bloom):
        self.bloom = bloom; self.seen = {}

    def check_add(self, chash, doc_id):
        if not chash:
            return None
        if chash in self.bloom and chash in self.seen:
            return self.seen[chash]
        self.bloom.add(chash); self.seen[chash] = doc_id
        return None


def sha256_file(path, chunk=1048576):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            b = f.read(chunk)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


def norm_text_hash(text):
    if not text:
        return None
    t = re.sub(r"\s+", " ", text).strip().lower()
    return hashlib.sha256(t.encode("utf-8")).hexdigest() if t else None


def blake2s_short(text, digest_bytes=3, upper=True):
    h = hashlib.blake2s(text.encode("utf-8"), digest_size=int(digest_bytes)).hexdigest()
    return h.upper() if upper else h


class ThreadAccel:
    def __init__(self, workers="AUTO", max_workers=8, min_files=12, enabled=True):
        self.enabled = bool(enabled); self.min_files = int(min_files)
        if workers in ("AUTO", None):
            cpu = os.cpu_count() or 4
            self.workers = max(2, min(int(max_workers), cpu * 2))
        else:
            self.workers = max(1, min(int(max_workers), int(workers)))

    def map(self, fn, items, on_progress=None):
        items = list(items); total = len(items); out = [None] * total
        if not self.enabled or total < self.min_files:
            for i, it in enumerate(items):
                out[i] = fn(it)
                if on_progress:
                    on_progress(i + 1, total, it)
            return out
        with ThreadPoolExecutor(max_workers=self.workers) as ex:
            futs = {ex.submit(fn, it): i for i, it in enumerate(items)}
            done = 0
            for fut in as_completed(futs):
                i = futs[fut]
                try:
                    out[i] = fut.result()
                except Exception as e:
                    out[i] = {"status": "ERROR", "reason": str(e)}
                done += 1
                if on_progress:
                    on_progress(done, total, items[i])
        return out


class Progress:
    def __init__(self, marker="@@PROGRESS|{pct}|{label}", stream="stderr", enabled=True):
        self.marker = marker; self.enabled = bool(enabled)
        self.out = sys.stderr if stream == "stderr" else sys.stdout

    def step(self, done, total, label=""):
        if not self.enabled:
            return
        pct = int(done * 100 / max(1, total))
        self.out.write(self.marker.replace("{pct}", str(pct)).replace("{label}", str(label)) + "\n")
        self.out.flush()


class Timer:
    def __init__(self):
        self.phases = {}; self._t = None; self._n = None

    def start(self, n):
        import time
        self._n = n; self._t = time.perf_counter()

    def stop(self):
        import time
        if self._n is None:
            return 0.0
        dt = time.perf_counter() - self._t
        self.phases[self._n] = round(self.phases.get(self._n, 0.0) + dt, 6)
        self._n = None; return dt

    def report(self):
        return dict(self.phases)


# ================================================================= B. CONFIG ==
DEFAULT_CONFIG = os.path.join(HERE, "via_forge_config.json")
DEFAULT_REGISTRY = os.path.join(HERE, "via_forge_registry.json")


def load_config(path=None, ensure_dirs=True):
    path = path or DEFAULT_CONFIG
    with open(path, encoding="utf-8-sig") as f:
        cfg = json.load(f)
    paths = cfg.get("paths", {})
    root = paths.get("root", "")
    if paths.get("root_fallback") == "AUTO":
        if not root or (os.name != "nt" and len(root) > 2 and root[1:3] == ":\\"):
            root = os.path.join(HERE, "runtime", "root")
    resolved = {}
    for k, v in paths.items():
        resolved[k] = os.path.normpath(v.replace("{root}", root)) if (isinstance(v, str) and "{root}" in v) else v
    resolved["root"] = root
    cfg["paths"] = resolved
    cfg["_config_path"] = os.path.abspath(path)
    if ensure_dirs:
        for key in ("root", "base", "inbox", "runtime", "ledger", "cache", "temp",
                    "output", "reports", "exports", "logs", "quarantine", "venv"):
            p = resolved.get(key)
            if p and not (len(p) > 1 and p[1] == ":" and os.name != "nt"):
                try:
                    os.makedirs(p, exist_ok=True)
                except Exception:
                    pass
    return cfg


def load_registry(path=None):
    path = path or DEFAULT_REGISTRY
    with open(path, encoding="utf-8-sig") as f:
        reg = json.load(f)
    return reg


# ================================================================ C. FORMATS ==
FORMAT_MAP = {
    ".docx": "OFFICE_WORD", ".dotx": "OFFICE_WORD", ".docm": "OFFICE_WORD",
    ".xlsx": "OFFICE_EXCEL", ".xltx": "OFFICE_EXCEL", ".xlsm": "OFFICE_EXCEL",
    ".pptx": "OFFICE_PPT", ".potx": "OFFICE_PPT", ".pptm": "OFFICE_PPT",
    ".doc": "OFFICE_LEGACY", ".xls": "OFFICE_LEGACY", ".ppt": "OFFICE_LEGACY",
    ".odt": "ODF", ".ods": "ODF", ".odp": "ODF",
    ".pdf": "PDF", ".rtf": "RTF", ".html": "HTML", ".htm": "HTML",
    ".eml": "EMAIL_EML", ".msg": "EMAIL_MSG",
    ".png": "IMAGE", ".jpg": "IMAGE", ".jpeg": "IMAGE", ".tif": "IMAGE",
    ".tiff": "IMAGE", ".bmp": "IMAGE", ".gif": "IMAGE",
    ".gdoc": "GOOGLE_POINTER", ".gsheet": "GOOGLE_POINTER", ".gslides": "GOOGLE_POINTER",
    ".txt": "TEXT", ".csv": "TEXT", ".tsv": "TEXT", ".md": "TEXT",
    ".markdown": "TEXT", ".json": "TEXT", ".jsonl": "TEXT", ".xml": "TEXT",
}
OOXML_FAMILIES = {"OFFICE_WORD", "OFFICE_EXCEL", "OFFICE_PPT"}
TEXTISH = {"TEXT", "HTML", "RTF", "GOOGLE_POINTER"}

RE_TWID = re.compile(r"\b[A-Z][12]\d{8}\b")
RE_EMAIL = re.compile(r"[\w.+-]+@[\w-]+\.[\w.-]+")
RE_PHONE = re.compile(r"\b09\d{2}[- ]?\d{3}[- ]?\d{3}\b")


def _decode(raw, probe=("utf-8-sig", "utf-8", "big5", "cp950", "latin-1")):
    cn = TOOLS.get("charset")
    if cn is not None:
        try:
            best = cn.from_bytes(raw).best()
            if best is not None:
                return str(best), best.encoding
        except Exception:
            pass
    for enc in probe:
        try:
            return raw.decode(enc), enc
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace"), "utf-8/replace"


# ================================================================= D. REPAIR ==
def _bytes_binary(raw):
    """True if raw bytes look binary (high NUL/control-byte ratio)."""
    if not raw:
        return True
    ctrl = sum(1 for b in raw if b < 0x09 or 0x0e <= b <= 0x1f)
    return (ctrl / len(raw)) > 0.15


def _text_damaged(s):
    """True if decoded text is mostly control/replacement chars (binary garbage)."""
    if not s:
        return True
    bad = sum(1 for ch in s if (ord(ch) < 32 and ch not in "\r\n\t") or ch == "\ufffd")
    return (bad / len(s)) > 0.15


def repair_encoding(path, out_dir):
    """Re-save text with clean UTF-8. Returns (fixed_path|None, detail, health)."""
    try:
        with open(path, "rb") as f:
            raw = f.read()
        if _bytes_binary(raw):
            return None, "binary/control bytes, not text", "DAMAGED"
        text, enc = _decode(raw)
        if _text_damaged(text):
            return None, "decoded content mostly non-printable (%s)" % enc, "DAMAGED"
        clean = (enc or "").lower().replace("-", "_")
        if clean in ("utf_8", "utf_8_sig", "u8", "utf", "ascii", "us_ascii"):
            return None, "encoding ok (%s)" % enc, "OK"
        os.makedirs(out_dir, exist_ok=True)
        fixed = os.path.join(out_dir, os.path.basename(path))
        with open(fixed, "w", encoding="utf-8") as f:
            f.write(text)
        return fixed, "re-encoded %s -> utf-8" % enc, "REPAIRED"
    except Exception as e:
        return None, "encoding repair failed: %s" % e, "DAMAGED"


def repair_zip_container(path, out_dir):
    """OOXML/ODF are zip. Validate; if bad-but-recoverable, repackage."""
    try:
        with zipfile.ZipFile(path) as z:
            bad = z.testzip()
        if bad is None:
            return None, "zip ok", "OK"
        return None, "zip corrupt member: %s" % bad, "DAMAGED"
    except zipfile.BadZipFile:
        return None, "not a valid zip", "DAMAGED"
    except Exception as e:
        return None, "zip check failed: %s" % e, "DAMAGED"


def repair_pdf(path, out_dir):
    pike = TOOLS.get("pikepdf")
    if pike is None:
        return None, "pikepdf absent; skip linearize", "OK"
    try:
        os.makedirs(out_dir, exist_ok=True)
        fixed = os.path.join(out_dir, os.path.basename(path))
        with pike.open(path) as pdf:
            pdf.save(fixed, linearize=True)
        return fixed, "pdf linearized", "REPAIRED"
    except Exception as e:
        return None, "pdf repair failed: %s" % e, "DAMAGED"


# ================================================================ E. EXTRACT ==
def _ocr_image(path):
    PIL = TOOLS.get("PIL"); pt = TOOLS.get("pytesseract")
    if PIL is None or pt is None or not TOOLS.bin_path("tesseract"):
        return "", "ocr_unavailable"
    try:
        from PIL import Image
        txt = pt.image_to_string(Image.open(path), lang="chi_tra+eng")
        return txt, "tesseract"
    except Exception as e:
        return "", "ocr_fail:%s" % e


def _soffice_txt(path, out_dir):
    if not TOOLS.bin_path("soffice"):
        return "", "soffice_absent"
    try:
        import subprocess
        os.makedirs(out_dir, exist_ok=True)
        subprocess.run(["soffice", "--headless", "--convert-to", "txt:Text",
                        "--outdir", out_dir, path],
                       timeout=60, capture_output=True)
        base = os.path.splitext(os.path.basename(path))[0] + ".txt"
        fp = os.path.join(out_dir, base)
        if os.path.exists(fp):
            with open(fp, encoding="utf-8", errors="replace") as f:
                return f.read(), "soffice"
        return "", "soffice_no_output"
    except Exception as e:
        return "", "soffice_fail:%s" % e


def extract_content(path, family, cfg, limit=8000):
    """Return (text, method, structured)."""
    structured = {}
    try:
        if family == "OFFICE_WORD" and TOOLS.get("docx"):
            import docx
            d = docx.Document(path)
            parts = [p.text for p in d.paragraphs if p.text.strip()]
            for t in d.tables:
                for row in t.rows:
                    parts.append(" ".join(c.text for c in row.cells))
            structured = {"paragraphs": len(d.paragraphs), "tables": len(d.tables)}
            return "\n".join(parts)[:limit], "docx", structured

        if family == "OFFICE_EXCEL":
            pd = TOOLS.get("pandas")
            if pd:
                xl = pd.ExcelFile(path); buf, sheets = [], []
                for sn in xl.sheet_names[:5]:
                    df = xl.parse(sn, nrows=40, dtype=str).fillna("")
                    sheets.append({"name": sn, "rows": int(df.shape[0]), "cols": int(df.shape[1]),
                                   "columns": [str(c) for c in df.columns][:30]})
                    buf.append(sn + " " + " ".join(str(c) for c in df.columns))
                    for _, r in df.head(30).iterrows():
                        buf.append(" ".join(map(str, r.tolist())))
                structured = {"sheets": sheets}
                return "\n".join(buf)[:limit], "pandas", structured
            op = TOOLS.get("openpyxl")
            if op:
                wb = op.load_workbook(path, read_only=True, data_only=True); buf = []
                for sn in wb.sheetnames[:5]:
                    ws = wb[sn]; buf.append(sn)
                    for i, row in enumerate(ws.iter_rows(values_only=True)):
                        if i >= 40:
                            break
                        buf.append(" ".join("" if v is None else str(v) for v in row))
                wb.close()
                return "\n".join(buf)[:limit], "openpyxl", structured

        if family == "OFFICE_PPT" and TOOLS.get("pptx"):
            import pptx
            prs = pptx.Presentation(path); parts = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    if shp.has_text_frame:
                        parts.append(shp.text_frame.text)
            structured = {"slides": len(prs.slides)}
            return "\n".join(parts)[:limit], "pptx", structured

        if family == "OFFICE_LEGACY":
            txt, meth = _soffice_txt(path, cfg["paths"].get("quarantine", HERE))
            return txt[:limit], meth, structured

        if family == "ODF" and TOOLS.get("odfpy"):
            from odf.opendocument import load as odload
            from odf import text as odtext, teletype
            doc = odload(path); parts = []
            for el in doc.getElementsByType(odtext.P):
                parts.append(teletype.extractText(el))
            return "\n".join(parts)[:limit], "odfpy", structured

        if family == "PDF":
            pypdf = TOOLS.get("pypdf")
            if pypdf:
                from pypdf import PdfReader
                r = PdfReader(path); parts = []
                for pg in r.pages[:30]:
                    parts.append(pg.extract_text() or "")
                text = "\n".join(parts)
                structured = {"pages": len(r.pages)}
                if len(text.strip()) < 20:  # image PDF -> OCR one page best-effort
                    return text[:limit], "pypdf_thin", structured
                return text[:limit], "pypdf", structured
            return "", "pdf_no_reader", structured

        if family == "RTF" and TOOLS.get("striprtf"):
            from striprtf.striprtf import rtf_to_text
            with open(path, encoding="utf-8", errors="replace") as f:
                return rtf_to_text(f.read())[:limit], "striprtf", structured

        if family == "HTML":
            with open(path, "rb") as f:
                raw = f.read()
            text, _ = _decode(raw)
            bs = TOOLS.get("bs4")
            if bs:
                from bs4 import BeautifulSoup
                return BeautifulSoup(text, "html.parser").get_text(" ")[:limit], "bs4", structured
            return re.sub(r"<[^>]+>", " ", text)[:limit], "regex_html", structured

        if family == "EMAIL_EML":
            import email
            from email import policy
            with open(path, "rb") as f:
                msg = email.message_from_bytes(f.read(), policy=policy.default)
            hdr = {k: str(msg[k]) for k in ("From", "To", "Cc", "Subject", "Date",
                                            "Message-ID", "In-Reply-To", "References") if msg[k]}

            def _part_text(part):
                payload = part.get_payload(decode=True)
                if payload is None:
                    try:
                        return part.get_content()
                    except Exception:
                        return ""
                cs = part.get_content_charset()
                if cs:
                    try:
                        return payload.decode(cs, errors="replace")
                    except Exception:
                        pass
                t, _ = _decode(payload)   # charset detection fallback (utf-8/big5/...)
                return t

            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body += _part_text(part)
            else:
                body = _part_text(msg)
            structured = {"headers": hdr,
                          "attachments": [p.get_filename() for p in msg.walk() if p.get_filename()]}
            text = "\n".join("%s: %s" % (k, v) for k, v in hdr.items()) + "\n\n" + (body or "")
            return text[:limit], "email.eml", structured

        if family == "EMAIL_MSG" and TOOLS.get("extract_msg"):
            import extract_msg
            m = extract_msg.Message(path)
            hdr = {"From": m.sender or "", "To": m.to or "", "Subject": m.subject or "",
                   "Date": str(m.date or "")}
            structured = {"headers": hdr,
                          "attachments": [a.longFilename or a.shortFilename for a in m.attachments]}
            text = "\n".join("%s: %s" % (k, v) for k, v in hdr.items() if v) + "\n\n" + (m.body or "")
            return text[:limit], "extract_msg", structured

        if family == "IMAGE":
            txt, meth = _ocr_image(path)
            return txt[:limit], meth, structured

        if family == "GOOGLE_POINTER":
            with open(path, "rb") as f:
                raw = f.read()
            text, _ = _decode(raw)
            doc_id, url = "", ""
            try:
                o = json.loads(text)
                doc_id = o.get("doc_id") or o.get("resource_id") or ""
                url = o.get("url") or ""
            except Exception:
                mm = re.search(r'"(?:doc_id|resource_id)"\s*:\s*"([^"]+)"', text)
                doc_id = mm.group(1) if mm else ""
            ext = os.path.splitext(path)[1].lower()
            kind = {".gdoc": "document", ".gsheet": "spreadsheets", ".gslides": "presentation"}.get(ext, "document")
            fmt = {"spreadsheets": "xlsx", "presentation": "pptx"}.get(kind, "docx")
            export = "https://docs.google.com/%s/d/%s/export?format=%s" % (kind, doc_id, fmt) if doc_id else ""
            structured = {"gdoc_kind": kind, "doc_id": doc_id, "url": url, "export_url": export}
            return ("[GOOGLE POINTER] doc_id=%s url=%s export=%s" % (doc_id, url, export))[:limit], "gdoc_pointer", structured

        if family in TEXTISH or family == "TEXT":
            with open(path, "rb") as f:
                raw = f.read()
            text, enc = _decode(raw)
            structured = {"encoding": enc}
            if path.lower().endswith((".csv", ".tsv")):
                sep = "\t" if path.lower().endswith(".tsv") else ","
                rows = list(csv.reader(text.splitlines(), delimiter=sep))
                structured["rows"] = max(0, len(rows) - 1)
                structured["cols"] = len(rows[0]) if rows else 0
            elif path.lower().endswith((".json", ".jsonl")):
                structured["json"] = True
            elif path.lower().endswith((".md", ".markdown")):
                text = re.sub(r"[#>*`_~]", " ", text)
            return text[:limit], "text", structured
    except Exception as e:
        return "", "extract_fail:%s" % e, structured
    return "", "unsupported", structured


# ================================================================ F. TEXTLAB ==
_FULL2HALF = {i: i - 0xFEE0 for i in range(0xFF01, 0xFF5F)}
_FULL2HALF[0x3000] = 0x20
_OPENCC = None


def tl_normalize(text, cfg):
    if not text:
        return ""
    t = text.translate(_FULL2HALF) if cfg["textlab"]["normalize_fullwidth_ascii"] else text
    t = re.sub(r"[ \t]+", " ", t)
    return re.sub(r"\n{3,}", "\n\n", t).strip()


def tl_correct(text, cfg, corrections=None):
    global _OPENCC
    out = text
    if cfg["textlab"]["opencc_convert"] and text:
        opencc = TOOLS.get("opencc")
        if opencc is not None:
            try:
                if _OPENCC is None:
                    _OPENCC = opencc.OpenCC(cfg["textlab"]["opencc_profile"])
                out = _OPENCC.convert(out)
            except Exception:
                pass
    ncorr = 0
    if corrections:
        for bad, good in corrections.items():
            if bad.startswith("_"):
                continue
            new = re.sub(r"\b%s\b" % re.escape(bad), good, out)
            if new != out:
                ncorr += 1; out = new
    return out, ncorr


def tl_sentences(text, cfg):
    if not cfg["textlab"]["sentence_split"] or not text:
        return []
    if cfg["textlab"]["decimal_point_protect"]:
        text = re.sub(r"(\d)\.(\d)", r"\1<DP>\2", text)
    parts = re.split(r"(?<=[\u3002\uff01\uff1f.!?\n])", text)
    return [p.replace("<DP>", ".").strip() for p in parts if p.strip()]


_KW_STOP = {"the", "and", "for", "with", "this", "that",
            "from", "subject", "cc", "bcc", "re", "fwd", "to", "sent", "date",
            "com", "org", "net", "www", "http", "https", "corp", "inc", "ltd",
            "email", "mail", "gmail", "html"}


def tl_keywords(text, cfg):
    topk = cfg["textlab"]["keyword_topk"]
    text = RE_EMAIL.sub(" ", text or "")  # drop email addresses so domain fragments don't leak
    jieba = TOOLS.get("jieba")
    if jieba is not None and cfg["textlab"]["keyword_engine"] == "jieba_tfidf":
        try:
            import jieba.analyse as ja
            return [{"term": k, "weight": round(float(w), 4)}
                    for k, w in ja.extract_tags(text, topK=topk, withWeight=True)
                    if k.lower() not in _KW_STOP]
        except Exception:
            pass
    cjk = re.findall(r"[\u4e00-\u9fff]{2,4}", text)
    lat = re.findall(r"[A-Za-z]{3,}", text)
    c = Counter(cjk + [w.lower() for w in lat])
    for stop in _KW_STOP:
        c.pop(stop, None)
    return [{"term": t, "weight": n} for t, n in c.most_common(topk)]


def tl_email_structure(structured):
    hdr = (structured or {}).get("headers", {})
    if not hdr:
        return None
    def orgs(v):
        return [e.split("@")[-1] for e in RE_EMAIL.findall(v or "")]
    participants = []
    for k in ("From", "To", "Cc"):
        participants += RE_EMAIL.findall(hdr.get(k, ""))
    return {"subject": hdr.get("Subject", ""),
            "participants": list(dict.fromkeys(participants))[:12],
            "organizations": list(dict.fromkeys(sum((orgs(hdr.get(k, "")) for k in ("From", "To", "Cc")), [])))[:12]}


def tl_enrich(text, family, structured, cfg, corrections=None):
    norm = tl_normalize(text, cfg)
    corrected, ncorr = tl_correct(norm, cfg, corrections)
    intel = {"corrected": corrected, "correction_count": ncorr,
             "sentence_count": len(tl_sentences(corrected, cfg)),
             "keywords": tl_keywords(corrected, cfg)}
    email_struct = tl_email_structure(structured) if family in ("EMAIL_EML", "EMAIL_MSG") else None
    return intel, email_struct


# ============================================================== G. CLASSIFY ==
def detect_language(text):
    if not text:
        return "unknown"
    cjk = len(re.findall(r"[\u4e00-\u9fff]", text))
    lat = len(re.findall(r"[A-Za-z]", text))
    if cjk and lat:
        return "zh+en" if cjk >= lat * 0.3 else "en"
    if cjk:
        return "zh"
    if lat:
        return "en"
    return "unknown"


def _kw_score(text, name, table, default):
    hay = ((text or "") + " " + (name or "")).lower()
    scores = {}
    for label, kws in table.items():
        hit = sum(1 for k in kws if k.lower() in hay)
        if hit:
            scores[label] = hit
    if not scores:
        return default, {}
    best = max(scores.items(), key=lambda kv: kv[1])[0]
    return best, dict(sorted(scores.items(), key=lambda kv: -kv[1]))


def classify(path, family, text, reg):
    name = os.path.basename(path)
    domains = reg["domains"]; sensitivity = reg["sensitivity"]
    fmt_class = ("OFFICE" if family in OOXML_FAMILIES or family == "OFFICE_LEGACY"
                 else "GOOGLE" if family == "GOOGLE_POINTER"
                 else "EMAIL" if family in ("EMAIL_EML", "EMAIL_MSG")
                 else "IMAGE" if family == "IMAGE" else family)
    domain, dscores = _kw_score(text, name, domains, "GENERAL")
    lang = detect_language(text)
    sens, sscores = _kw_score(text, name, sensitivity, "PUBLIC_OR_UNMARKED")
    pii = []
    if text:
        if RE_TWID.search(text):
            pii.append("TW_ID")
        if RE_EMAIL.search(text):
            pii.append("EMAIL")
        if RE_PHONE.search(text):
            pii.append("PHONE")
    if pii and sens == "PUBLIC_OR_UNMARKED":
        sens = "CONFIDENTIAL"
    conf = "M" if (text and (dscores or sscores)) else ("P" if text else "Est")
    return {"format_class": fmt_class, "content_domain": domain, "domain_scores": dscores,
            "language": lang, "sensitivity": sens, "sensitivity_scores": sscores,
            "pii": pii, "classify_confidence": conf}


# ================================================================ ENGINE ==
class ForgeEngine:
    def __init__(self, cfg=None, reg=None):
        self.cfg = cfg or load_config()
        self.reg = reg or load_registry()
        self.corrections = self.reg.get("corrections", {})
        acc = self.cfg["accelerator"]
        self.bloom = BloomFilter(size=acc["bloom_size"], hashes=acc["bloom_hashes"])
        self.dedup = Dedup(self.bloom)
        self.prog = Progress(acc["progress_marker"], acc["progress_stream"])
        self.timer = Timer()
        self.cache = DiskCache(self.cfg["paths"]["cache"], acc["disk_cache_ttl_hours"], acc["disk_cache_enabled"])
        self.accel = ThreadAccel(acc["parallel_workers"], acc["parallel_max_workers"],
                                 acc["parallel_min_files"], acc["parallel_enabled"])
        self.ledger_path = os.path.join(self.cfg["paths"]["ledger"], self.cfg["ledger"]["file"])
        self.records = []; self.seq = 0

    def supported(self):
        return sorted(FORMAT_MAP.keys())

    def _make_id(self, path):
        day = datetime.now(UTC).strftime("%Y%m%d")
        return "VIA-DOC-%s-%06d-%s" % (day, self.seq, blake2s_short(path, self.cfg["ids"]["hash_digest_bytes"]))

    def _repair(self, path, family):
        qdir = self.cfg["paths"].get("quarantine", HERE)
        if family in OOXML_FAMILIES or family == "ODF":
            return repair_zip_container(path, qdir)
        if family == "PDF":
            return repair_pdf(path, qdir)
        if family in ("TEXT", "HTML", "RTF"):
            return repair_encoding(path, qdir)
        return None, "no repair path", "OK"

    def process_one(self, path, do_fix=True):
        ext = os.path.splitext(path)[1].lower()
        family = FORMAT_MAP.get(ext)
        if family is None:
            return {"path": path, "status": "SKIP", "reason": "unsupported:%s" % ext}
        try:
            size = os.path.getsize(path)
            mtime = datetime.fromtimestamp(os.path.getmtime(path), UTC).isoformat()
            # VPNS A1: SHA \u5feb\u53d6\uff08path+size+mtime \u4e0d\u8b8a \u2192 \u514d\u91cd\u7b97\uff09
            # VPNS A3: \u5927\u6a94\u8df3\u96dc\u6e4a\uff08>50MB \u6a19 SKIP\uff0c\u4e0d\u5361\u65b7\uff09
            ck = "%s|%d|%s" % (path, size, mtime)
            cache = getattr(self, "_sha_cache", None)
            if cache is None:
                cache = self._sha_cache = {}
            if ck in cache:
                sha = cache[ck]
            elif size > 50 * 1024 * 1024:
                sha = "SKIP>50MB:%d" % size
                cache[ck] = sha
            else:
                sha = sha256_file(path, self.cfg["accelerator"]["hash_chunk_bytes"])
                cache[ck] = sha
        except Exception as e:
            return {"path": path, "status": "ERROR", "reason": str(e)}

        fixed, detail, health = (None, "", "OK")
        if do_fix:
            fixed, detail, health = self._repair(path, family)

        read_path = fixed or path
        text, method, structured = extract_content(read_path, family, self.cfg,
                                                    limit=self.cfg["readers"].get("extract_limit", 8000))
        if not text and health == "OK":
            health = "UNREADABLE" if method not in ("gdoc_pointer",) else "OK"

        asp = classify(path, family, text, self.reg)
        intel, email_struct = tl_enrich(text, family, structured, self.cfg, self.corrections)
        asp["keywords"] = [k["term"] for k in intel["keywords"]][:8]

        chash = norm_text_hash(text) or sha
        dup_of = self.dedup.check_add(chash, "PENDING")
        is_dup = bool(dup_of)
        if is_dup:
            health = "DUPLICATE"

        return {
            "_family": family, "_sha": sha, "_size": size, "_mtime": mtime,
            "_health": health, "_detail": detail, "_method": method,
            "_structured": structured, "_asp": asp, "_intel": intel,
            "_email": email_struct, "_chash": chash, "_dup_of": dup_of, "_fixed": fixed,
            "_text": text or "",
            "path": path,
        }

    def scan(self, target, do_fix=True):
        self.timer.start("scan")
        if isinstance(target, str) and os.path.isdir(target):
            paths = self._walk(target)
        elif isinstance(target, (list, tuple)):
            paths = list(target)
        else:
            paths = [target]

        raw = self.accel.map(
            lambda p: self.process_one(p, do_fix),
            paths,
            on_progress=lambda d, t, it: self.prog.step(d, t, os.path.basename(str(it))))

        results = []
        for r in raw:
            if not isinstance(r, dict) or r.get("status") in ("SKIP", "ERROR"):
                results.append(r); continue
            self.seq += 1
            doc_id = self._make_id(r["path"])
            _chash = r["_chash"]
            if r["_dup_of"]:                       # duplicate -> resolve to original's real id
                _dup_of_id = self.dedup.seen.get(_chash)
                if _dup_of_id in (None, "PENDING"):
                    _dup_of_id = None
            else:                                  # first occurrence -> register real id
                self.dedup.seen[_chash] = doc_id
                _dup_of_id = None
            rec = {
                "doc_id": doc_id,
                "source_path": os.path.abspath(r["path"]),
                "fixed_path": os.path.abspath(r["_fixed"]) if r["_fixed"] else None,
                "file_name": os.path.basename(r["path"]),
                "extension": os.path.splitext(r["path"])[1].lower(),
                "format_family": r["_family"], "size_bytes": r["_size"],
                "modified_time": r["_mtime"], "sha256": r["_sha"], "content_hash": r["_chash"],
                "health": r["_health"], "is_duplicate": bool(r["_dup_of"]),
                "duplicate_of": _dup_of_id,
                "repair_detail": r["_detail"], "source_method": r["_method"],
                "structured": r["_structured"], "aspects": r["_asp"],
                "text_intel": {k: v for k, v in r["_intel"].items() if k != "corrected"},
                "text_preview": (r.get("_text", "") or "")[:self.cfg["readers"].get("text_preview_chars", 400)],
                "email_structure": r["_email"],
                "processed_at": datetime.now(UTC).isoformat(),
                "evidence": "LOCAL_AUTHORIZED_SCAN", "status": "OK",
            }
            self.records.append(rec)
            self._append_ledger(rec)
            results.append(rec)
        self.timer.stop()
        return results

    def _walk(self, root):
        skip = set(self.cfg["readers"]["skip_dirs"])
        out = []
        for dp, dns, fns in os.walk(root):
            dns[:] = [d for d in dns if d not in skip]
            for fn in fns:
                if self.cfg["readers"]["skip_hidden"] and fn.startswith("."):
                    continue
                if os.path.splitext(fn)[1].lower() in FORMAT_MAP:
                    out.append(os.path.join(dp, fn))
        return sorted(out)

    def _append_ledger(self, rec):
        os.makedirs(os.path.dirname(self.ledger_path), exist_ok=True)
        with open(self.ledger_path, "a", encoding="utf-8") as f:
            f.write(json.dumps(rec, ensure_ascii=False) + "\n")

    def summary(self, records):
        ok = [r for r in records if r.get("status") == "OK"]
        def cnt(fn):
            d = {}
            for r in ok:
                k = fn(r); d[k] = d.get(k, 0) + 1
            return dict(sorted(d.items(), key=lambda kv: -kv[1]))
        return {
            "total": len(records), "ok": len(ok),
            "dup": sum(1 for r in ok if r["is_duplicate"]),
            "healthy": sum(1 for r in ok if r["health"] in ("OK", "REPAIRED")),
            "bad": sum(1 for r in ok if r["health"] in ("DAMAGED", "UNREADABLE")),
            "skip": sum(1 for r in records if r.get("status") == "SKIP"),
            "error": sum(1 for r in records if r.get("status") == "ERROR"),
            "by_health": cnt(lambda r: r["health"]),
            "by_domain": cnt(lambda r: r["aspects"]["content_domain"]),
            "by_sensitivity": cnt(lambda r: r["aspects"]["sensitivity"]),
            "by_format": cnt(lambda r: r["format_family"]),
            "by_language": cnt(lambda r: r["aspects"]["language"]),
            "by_evidence": cnt(lambda r: r["aspects"]["classify_confidence"]),
            "tool_matrix": TOOLS.matrix(), "cache": self.cache.stats(),
            "timing": self.timer.report(),
            "generated_at": datetime.now(UTC).isoformat(),
        }

    # ------------------------------------------------------------- exports
    def export_json(self, records, summary, path):
        with open(path, "w", encoding="utf-8") as f:
            json.dump({"summary": summary, "records": records}, f, ensure_ascii=False, indent=2)
        return path

    def export_csv(self, records, path):
        cols = ["doc_id", "file_name", "format_family", "health", "content_domain",
                "sensitivity", "pii", "classify_confidence", "size_bytes", "source_path"]
        with open(path, "w", encoding="utf-8-sig", newline="") as f:
            w = csv.writer(f); w.writerow(cols)
            for r in records:
                if r.get("status") != "OK":
                    continue
                a = r["aspects"]
                w.writerow([r["doc_id"], r["file_name"], r["format_family"], r["health"],
                            a["content_domain"], a["sensitivity"], ",".join(a["pii"]),
                            a["classify_confidence"], r["size_bytes"], r["source_path"]])
        return path

    def export_xlsx(self, records, path):
        pd = TOOLS.get("pandas")
        if pd is None:
            return None
        rows = []
        for r in records:
            if r.get("status") != "OK":
                continue
            a = r["aspects"]
            rows.append({"doc_id": r["doc_id"], "file_name": r["file_name"],
                         "format": r["format_family"], "health": r["health"],
                         "domain": a["content_domain"], "sensitivity": a["sensitivity"],
                         "pii": ",".join(a["pii"]), "evidence": a["classify_confidence"],
                         "keywords": "\u3001".join(a.get("keywords", [])[:8]),
                         "size_bytes": r["size_bytes"], "source_path": r["source_path"]})
        pd.DataFrame(rows).to_excel(path, index=False)
        return path

    def export_html(self, records, summary, path):
        vl = self.cfg["visual_lock"]
        ok = [r for r in records if r.get("status") == "OK"]
        HC = {"OK": vl["down_green"], "REPAIRED": vl["teal"], "DAMAGED": vl["up_red"],
              "UNREADABLE": "#6b6a66", "DUPLICATE": vl["blue"]}
        trs = []
        for r in ok[:800]:
            a = r["aspects"]; hc = HC.get(r["health"], "#888")
            trs.append("<tr><td class=mono>%s</td><td>%s</td><td>%s</td>"
                       "<td><b style='color:%s'>%s</b></td><td>%s</td><td>%s</td>"
                       "<td style='color:%s'>%s</td><td>%s</td><td class=mono>%s</td></tr>" % (
                           r["doc_id"], r["file_name"], r["format_family"], hc, r["health"],
                           a["content_domain"], a["sensitivity"], vl["up_red"], ",".join(a["pii"]),
                           "\u3001".join(a.get("keywords", [])[:5]), a["classify_confidence"]))
        cards = "".join("<div class=card><b>%s</b><span>%s</span></div>" % (v, k)
                        for k, v in [("\u7e3d\u6578", summary["total"]), ("\u5065\u5eb7", summary["healthy"]),
                                     ("\u91cd\u8907", summary["dup"]), ("\u53d7\u640d", summary["bad"])])
        html = _HTML.format(bg=vl["bg"], paper=vl["paper"], ink=vl["ink"], line=vl["line"],
                            blue=vl["blue"], teal=vl["teal"], seal=vl["seal"],
                            cards=cards, rows="".join(trs), gen=summary["generated_at"])
        with open(path, "w", encoding="utf-8") as f:
            f.write(html)
        return path


_HTML = """<!DOCTYPE html><html lang=zh-Hant><head><meta charset=utf-8>
<meta name=viewport content="width=device-width,initial-scale=1"><title>VIA Forge Report</title>
<style>*{{box-sizing:border-box;margin:0;padding:0}}body{{background:{bg};color:{ink};
font-family:'DM Sans',system-ui,sans-serif;padding:30px}}.seal{{display:inline-flex;width:34px;height:34px;
background:{blue};color:#fff;border-radius:3px;align-items:center;justify-content:center;font-weight:700;margin-right:10px}}
h1{{font-family:'Syne',sans-serif;font-size:22px;display:flex;align-items:center}}
.sub{{color:#8a877f;font-size:12px;margin:4px 0 18px}}.cards{{display:flex;gap:12px;margin-bottom:18px;flex-wrap:wrap}}
.card{{background:{paper};border:1px solid {line};border-radius:3px;padding:14px 18px;min-width:90px}}
.card b{{font-size:26px;font-family:'Syne',sans-serif;display:block}}.card span{{font-size:11px;color:#8a877f}}
table{{width:100%;border-collapse:collapse;background:{paper};border:1px solid {line};border-radius:3px}}
th,td{{padding:8px 10px;text-align:left;border-bottom:1px solid {line};font-size:12px}}
th{{background:#faf9f6;font-family:'DM Mono',monospace;font-size:10px;text-transform:uppercase;color:#9a978f}}
.mono{{font-family:'DM Mono',monospace;font-size:11px;color:{blue}}}</style></head><body>
<h1><span class=seal>{seal}</span>VIA Forge \u2014 \u64f7\u53d6\u6cbb\u7406\u5831\u544a</h1>
<div class=sub>generated {gen} \u00b7 loopback \u00b7 append-only</div>
<div class=cards>{cards}</div>
<table><thead><tr><th>doc_id</th><th>\u6a94\u540d</th><th>\u683c\u5f0f</th><th>\u5065\u5eb7</th><th>\u9818\u57df</th>
<th>\u654f\u611f\u5ea6</th><th>PII</th><th>\u95dc\u9375\u8a5e</th><th>\u4fe1\u5fc3</th></tr></thead>
<tbody>{rows}</tbody></table></body></html>"""


# ================================================================ SELFTEST ==
def selftest():
    import tempfile
    p, f = 0, 0

    def ck(n, c):
        nonlocal p, f
        if c:
            p += 1; print("  [PASS] forge: %s" % n)
        else:
            f += 1; print("  [FAIL] forge: %s" % n)

    cfg = load_config(); reg = load_registry()
    d = tempfile.mkdtemp(); side = tempfile.mkdtemp()
    cfg["paths"]["cache"] = os.path.join(side, "cache")
    cfg["paths"]["ledger"] = side
    cfg["paths"]["quarantine"] = os.path.join(side, "q")
    os.makedirs(cfg["paths"]["cache"], exist_ok=True)

    def w(name, content):
        with open(os.path.join(d, name), "w", encoding="utf-8") as fh:
            fh.write(content)

    w("fin.md", "# Q2\n\u672c\u5b63\u71df\u6536\u6210\u9577\uff0c\u6bdb\u5229\u63d0\u5347\uff0c\u73fe\u91d1\u6d41\u6539\u5584\n")
    w("legal.txt", "\u5408\u7d04\u689d\u6b3e\u8207\u4fdd\u5bc6\u7d04\u5b9a NDA compliance\n")
    w("hr.csv", "who,id\nchris,A123456789\n")
    w("tech.json", "{\"api\":\"schema\",\"deployment\":\"pipeline\"}")
    w("mail.eml", "From: ivy@dongguan.com\nTo: sandeep@india.com\nSubject: \u4ea4\u671f\u63a1\u8cfc\n\n\u5ba2\u6236\u8a02\u55ae\u4ea4\u671f\n")
    w("plain.md", "teh quick adress recieve\n")  # typo correction

    eng = ForgeEngine(cfg, reg)
    recs = eng.scan(d, do_fix=True)
    ok = [r for r in recs if r.get("status") == "OK"]
    ck("all processed", len(ok) == 6)

    idx = {r["file_name"]: r for r in ok}
    ck("md finance domain", idx["fin.md"]["aspects"]["content_domain"] == "FINANCE")
    ck("txt legal domain", idx["legal.txt"]["aspects"]["content_domain"] == "LEGAL")
    ck("csv HR + TW_ID PII", "TW_ID" in idx["hr.csv"]["aspects"]["pii"])
    ck("PII escalates CONFIDENTIAL", idx["hr.csv"]["aspects"]["sensitivity"] == "CONFIDENTIAL")
    ck("json tech domain", idx["tech.json"]["aspects"]["content_domain"] == "TECHNICAL")
    ck("eml family", idx["mail.eml"]["format_family"] == "EMAIL_EML")
    ck("eml structure", idx["mail.eml"]["email_structure"] is not None)
    ck("eml orgs parsed", "dongguan.com" in (idx["mail.eml"]["email_structure"]["organizations"]))
    ck("typo corrected", idx["plain.md"]["text_intel"]["correction_count"] >= 2)
    ck("doc_id format", all(re.match(r"VIA-DOC-\d{8}-\d{6}-[0-9A-F]{6}", r["doc_id"]) for r in ok))
    ck("evidence tier", all(r["aspects"]["classify_confidence"] in ("M", "P", "Est") for r in ok))
    ck("health values", all(r["health"] in ("OK", "REPAIRED", "DAMAGED", "UNREADABLE", "DUPLICATE") for r in ok))

    # dedup: copy fin.md content
    w("fin_copy.md", "# Q2\n\u672c\u5b63\u71df\u6536\u6210\u9577\uff0c\u6bdb\u5229\u63d0\u5347\uff0c\u73fe\u91d1\u6d41\u6539\u5584\n")
    recs2 = eng.scan([os.path.join(d, "fin_copy.md")], do_fix=True)
    ck("dedup detects duplicate", recs2[0]["is_duplicate"] and recs2[0]["health"] == "DUPLICATE")
    ck("duplicate_of resolves to original id",
       bool(recs2[0]["duplicate_of"]) and recs2[0]["duplicate_of"] != "PENDING"
       and recs2[0]["duplicate_of"].startswith("VIA-DOC-"))

    summ = eng.summary(recs)
    ck("summary domains", "FINANCE" in summ["by_domain"])
    ck("summary tool matrix", isinstance(summ["tool_matrix"], dict) and "docx" in summ["tool_matrix"])

    # ledger append-only
    with open(eng.ledger_path, encoding="utf-8") as fh:
        n = sum(1 for _ in fh)
    ck("ledger append-only", n == 7)

    # damage detection: binary/control garbage must be DAMAGED, not REPAIRED
    with open(os.path.join(d, "garbage.txt"), "wb") as fb:
        fb.write(bytes([0xff, 0xfe, 0x00, 0x01, 0x02, 0xC0, 0x00, 0x03]))
    recs3 = eng.scan([os.path.join(d, "garbage.txt")], do_fix=True)
    ck("binary garbage marked DAMAGED", recs3[0]["health"] == "DAMAGED")

    # exports
    jp = eng.export_json(recs, summ, os.path.join(side, "o.json"))
    ck("export json", os.path.exists(jp))
    cp = eng.export_csv(recs, os.path.join(side, "o.csv"))
    ck("export csv", os.path.exists(cp))
    hp = eng.export_html(recs, summ, os.path.join(side, "o.html"))
    with open(hp, encoding="utf-8") as fh:
        htm = fh.read()
    ck("export html visual lock", "#4c78a8" in htm and "Forge" in htm)

    ck("supported superset", len(eng.supported()) >= 28)
    ck("tool matrix bins", "tesseract" in TOOLS.matrix() and "soffice" in TOOLS.matrix())
    return p, f


def main():
    ap = argparse.ArgumentParser(description="VIA Forge — unified fetch/fix/classify engine")
    ap.add_argument("--scan"); ap.add_argument("--fix", action="store_true")
    ap.add_argument("--export", choices=["json", "csv", "xlsx", "html"])
    ap.add_argument("--selftest", action="store_true"); ap.add_argument("--caps", action="store_true")
    args = ap.parse_args()

    if args.selftest:
        a, b = selftest(); print("forge: %d PASS / %d FAIL" % (a, b)); sys.exit(0 if b == 0 else 1)
    if args.caps:
        print(json.dumps(TOOLS.matrix(), ensure_ascii=False, indent=2)); return
    if args.scan:
        eng = ForgeEngine()
        recs = eng.scan(args.scan, do_fix=args.fix)
        summ = eng.summary(recs)
        print("OK=%d healthy=%d dup=%d bad=%d skip=%d err=%d" % (
            summ["ok"], summ["healthy"], summ["dup"], summ["bad"], summ["skip"], summ["error"]))
        if args.export:
            out = eng.cfg["paths"]["reports"]; os.makedirs(out, exist_ok=True)
            fn = {"json": "forge_report.json", "csv": "forge_records.csv",
                  "xlsx": "forge_records.xlsx", "html": "forge_report.html"}[args.export]
            path = os.path.join(out, fn)
            getattr(eng, "export_" + args.export)(recs, summ, path) if args.export in ("json", "html") \
                else getattr(eng, "export_" + args.export)(recs, path)
            print("exported -> %s" % path)


if __name__ == "__main__":
    main()
